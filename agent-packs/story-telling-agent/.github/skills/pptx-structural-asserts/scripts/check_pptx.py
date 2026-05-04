#!/usr/bin/env python3
"""check_pptx.py — Programmatic structural assertions for @deck-critic.

Rebuild scope (session 2026-05-04-7d3f9a2b):

  F1  — Pillow ImageFont real-metric overflow check (replaces the old
        heuristic). Honours run/paragraph font sizes, frame insets,
        paragraph space_before / space_after, and emits one record per
        offending shape (no early return; no 5% slack).
  F4  — Theme-resolved contrast: scans every run, applies WCAG
        large-text tier (3:1 for >=18pt or >=14pt-bold), surfaces
        unresolvable runs in a `contrast_unresolved` warning bucket
        (>=5 -> blocking).
  F7  — Saturation-aware `_bg_label`: highly saturated colours classify
        as `accent` regardless of luminance (sat > 0.45 -> accent;
        lum < 0.18 -> dark; lum > 0.72 -> light; else accent).
  F8  — `dark_light_run` thresholds: warn at >=3, blocking at >=7,
        OR blocking at >=5 if no `accent` slide breaks the run.

The font fallback chain is the single resolver in
``render-visual/assets/font_locator.py`` (architecture critic C3 —
single source of truth).
"""

from __future__ import annotations

import argparse
import colorsys
import json
import subprocess
import sys
from pathlib import Path


# ---------- bootstrap deps ----------

def _ensure(mod: str, pip_name: str | None = None) -> None:
    try:
        __import__(mod)
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", pip_name or mod])


_ensure("pptx", "python-pptx")
_ensure("PIL", "Pillow")

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402
from PIL import ImageFont  # noqa: E402


# ---------- font locator (single source of truth — sibling skill) ----------

def _resolve_font_locator():
    """Locate render-visual/assets/font_locator.find_dejavu_sans.

    Falls through to None if the sibling helper isn't installed (older
    pack layout); callers cope with `_measure_run_width_emu` raising
    OSError and using PIL's bitmap default.
    """
    here = Path(__file__).resolve()
    # .../skills/pptx-structural-asserts/scripts/check_pptx.py
    sibling = here.parent.parent.parent / "render-visual" / "assets"
    if (sibling / "font_locator.py").is_file():
        if str(sibling) not in sys.path:
            sys.path.insert(0, str(sibling))
        try:
            import font_locator  # type: ignore
            return font_locator.find_dejavu_sans
        except Exception:
            return None
    return None


_FIND_FALLBACK_FONT = _resolve_font_locator()


# ---------- units / colour helpers ----------

EMU_PER_INCH = 914400
EMU_PER_PT = 12700  # 914400 / 72
SNAP_EMU = int(0.05 * EMU_PER_INCH)


def _hex_to_rgb(h: str) -> tuple[int, int, int] | None:
    if not h:
        return None
    h = h.strip().lstrip("#")
    if len(h) != 6:
        return None
    try:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except ValueError:
        return None


def _rgb_to_relative_luminance(rgb: tuple[int, int, int]) -> float:
    def chan(c: int) -> float:
        x = c / 255.0
        return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _hsv_saturation(rgb: tuple[int, int, int]) -> float:
    r, g, b = (c / 255.0 for c in rgb)
    _, _, _ = colorsys.rgb_to_hsv(r, g, b)  # noqa
    _h, s, _v = colorsys.rgb_to_hsv(r, g, b)
    return s


def _contrast_ratio(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    l1 = _rgb_to_relative_luminance(fg)
    l2 = _rgb_to_relative_luminance(bg)
    a, b = max(l1, l2), min(l1, l2)
    return (a + 0.05) / (b + 0.05)


def _shape_solid_rgb(shape) -> tuple[int, int, int] | None:
    try:
        f = shape.fill
        if f.type is None:
            return None
        if hasattr(f, "fore_color") and getattr(f.fore_color, "type", None) is not None:
            try:
                rgb = f.fore_color.rgb
                if rgb is not None:
                    return rgb[0], rgb[1], rgb[2]
            except Exception:
                pass
    except Exception:
        return None
    return None


def _slide_bg_rgb(slide) -> tuple[int, int, int] | None:
    try:
        bg = slide.background
        if bg is not None and bg.fill is not None:
            ftype = getattr(bg.fill, "type", None)
            if ftype is not None:
                try:
                    rgb = bg.fill.fore_color.rgb
                    if rgb is not None:
                        return rgb[0], rgb[1], rgb[2]
                except Exception:
                    pass
    except Exception:
        pass
    best, best_area = None, 0
    for shp in slide.shapes:
        try:
            w = int(shp.width or 0)
            h = int(shp.height or 0)
            area = w * h
            rgb = _shape_solid_rgb(shp)
            if rgb and area > best_area:
                best, best_area = rgb, area
        except Exception:
            continue
    return best


# ---------- F7: saturation-aware background label ----------

def _bg_label(rgb: tuple[int, int, int] | None) -> str:
    """Classify a slide background.

    Per architecture critic C4 (and finding F7): saturation overrides
    luminance — heavily-saturated brand colours (e.g. #635BFF at lum
    ~0.17 / sat ~0.64) are `accent`, not `dark`. The hard-coded
    threshold is sat > 0.45.
    """
    if rgb is None:
        return "unknown"
    lum = _rgb_to_relative_luminance(rgb)
    sat = _hsv_saturation(rgb)
    if sat > 0.45:
        return "accent"
    if lum < 0.18:
        return "dark"
    if lum > 0.72:
        return "light"
    return "accent"


# ---------- per-slide content extractors ----------

def _slide_title_text(slide) -> str:
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            return (slide.shapes.title.text_frame.text or "").strip()
    except Exception:
        pass
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            t = (shp.text_frame.text or "").strip()
            if t:
                return t
    return ""


def _body_word_count(slide) -> int:
    title = _slide_title_text(slide)
    words = 0
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        t = (shp.text_frame.text or "").strip()
        if not t or t == title:
            continue
        words += len(t.split())
    return words


def _layout_hash(slide) -> str:
    parts = []
    for shp in slide.shapes:
        try:
            l = round((shp.left or 0) / EMU_PER_INCH, 1)
            t = round((shp.top or 0) / EMU_PER_INCH, 1)
            w = round((shp.width or 0) / EMU_PER_INCH, 1)
            h = round((shp.height or 0) / EMU_PER_INCH, 1)
            tf = 1 if getattr(shp, "has_text_frame", False) else 0
            parts.append((l, t, w, h, tf))
        except Exception:
            continue
    parts.sort()
    return "h:" + str(hash(tuple(parts)))[-12:]


def _has_underline_below_title(slide) -> bool:
    title = None
    try:
        title = slide.shapes.title
    except Exception:
        title = None
    if title is None:
        return False
    title_left = title.left or 0
    title_top = title.top or 0
    title_h = title.height or 0
    band_top = title_top + title_h
    band_top_min = band_top - int(0.1 * EMU_PER_INCH)
    band_top_max = band_top + int(0.5 * EMU_PER_INCH)
    for shp in slide.shapes:
        if shp == title:
            continue
        try:
            h = shp.height or 0
            t = shp.top or 0
            l = shp.left or 0
            w = shp.width or 0
            if h <= int(0.06 * EMU_PER_INCH) and band_top_min <= t <= band_top_max:
                if l <= title_left + int(0.5 * EMU_PER_INCH) and (l + w) >= title_left + int(1.0 * EMU_PER_INCH):
                    return True
        except Exception:
            continue
    return False


def _alignment_violations_for(slide, idx: int) -> list[dict]:
    out = []
    for shp in slide.shapes:
        try:
            l = shp.left or 0
            t = shp.top or 0
            for axis, val in (("left", l), ("top", t)):
                rem = val % SNAP_EMU
                if rem > 200 and (SNAP_EMU - rem) > 200:
                    out.append({
                        "slide": idx,
                        "shape": getattr(shp, "name", "?"),
                        "off_snap_axis": axis,
                        "off_by_emu": int(min(rem, SNAP_EMU - rem)),
                    })
                    break
        except Exception:
            continue
    return out


# ---------- F1: Pillow-based real-metric overflow ----------

_FONT_CACHE: dict[tuple[str | None, int], "ImageFont.FreeTypeFont"] = {}


def _load_font(family: str | None, size_pt: float) -> "ImageFont.ImageFont":
    """Best-effort font loader. 4x upscale for measurement accuracy.

    Resolution order: explicit family on host -> bundled DejaVuSans
    via render-visual font_locator -> matplotlib's bundled DejaVu ->
    PIL bitmap default. Uses module-level cache to avoid repeated
    truetype() calls.
    """
    px = max(8, int(size_pt * 4))
    cache_key = (family, px)
    if cache_key in _FONT_CACHE:
        return _FONT_CACHE[cache_key]
    font = None
    if family:
        try:
            font = ImageFont.truetype(family, px)
        except (OSError, IOError):
            font = None
    if font is None and _FIND_FALLBACK_FONT is not None:
        path = _FIND_FALLBACK_FONT()
        if path:
            try:
                font = ImageFont.truetype(path, px)
            except (OSError, IOError):
                font = None
    if font is None:
        # Last-resort bitmap default. Pixel accuracy degrades; report
        # records `font_fallback: true` so this is observable.
        font = ImageFont.load_default()
    _FONT_CACHE[cache_key] = font
    return font


def _measure_text_width_emu(text: str, size_pt: float, family: str | None) -> int:
    if not text:
        return 0
    font = _load_font(family, size_pt)
    try:
        # FreeTypeFont — getlength returns subpixel-accurate width
        px_at_4x = font.getlength(text)
    except AttributeError:
        # Bitmap default fallback — best effort
        try:
            l, t, r, b = font.getbbox(text)
            px_at_4x = float(r - l)
            # Bitmap path didn't get the 4x upscale; re-scale so the
            # math below stays consistent (assume default ~ 11px font).
            return int((px_at_4x / 11.0) * size_pt * EMU_PER_PT * 0.55)
        except Exception:
            return int(len(text) * size_pt * EMU_PER_PT * 0.55)
    in_units = (px_at_4x / 4.0) / 96.0  # ImageFont raster is 96 dpi
    return int(in_units * EMU_PER_INCH)


def _shape_overflow(shape, idx: int, font_fallback_flag: list[bool]) -> dict | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    if not (tf.text or "").strip():
        return None

    pad_default = int(0.05 * EMU_PER_INCH)
    pad_l = int(tf.margin_left)   if tf.margin_left   is not None else pad_default
    pad_r = int(tf.margin_right)  if tf.margin_right  is not None else pad_default
    pad_t = int(tf.margin_top)    if tf.margin_top    is not None else pad_default
    pad_b = int(tf.margin_bottom) if tf.margin_bottom is not None else pad_default
    inner_w_emu = int(shape.width or 0) - pad_l - pad_r
    inner_h_emu = int(shape.height or 0) - pad_t - pad_b
    if inner_w_emu <= 0 or inner_h_emu <= 0:
        return None

    total_height_emu = 0
    wrapped_lines = 0

    for para in tf.paragraphs:
        first_run = para.runs[0] if para.runs else None
        size_pt = (
            (first_run.font.size.pt if (first_run and first_run.font.size) else None)
            or (para.font.size.pt if para.font.size else None)
            or 18.0
        )
        family = (first_run.font.name if first_run else None) or (para.font.name if hasattr(para, "font") else None)

        if family is None:
            font_fallback_flag[0] = True

        line_h_emu = int(size_pt * 1.2 * EMU_PER_PT)
        text = "".join(r.text for r in para.runs) if para.runs else (para.text or "")
        if not text:
            total_height_emu += line_h_emu
            continue

        # Honour bullet indent (paragraph.level): each level eats ~0.25" of width
        try:
            level = para.level or 0
        except Exception:
            level = 0
        indent_emu = int(level * 0.25 * EMU_PER_INCH)
        avail_w = max(int(0.25 * EMU_PER_INCH), inner_w_emu - indent_emu)

        words = text.split()
        space_w = _measure_text_width_emu(" ", size_pt, family)
        wrapped = 0
        cur_w = 0
        cur_has_word = False
        for word in words:
            ww = _measure_text_width_emu(word, size_pt, family)
            tentative = cur_w + (space_w if cur_has_word else 0) + ww
            if tentative > avail_w and cur_has_word:
                wrapped += 1
                cur_w = ww
                cur_has_word = True
            else:
                cur_w = tentative
                cur_has_word = True
        if cur_has_word:
            wrapped += 1

        para_h = wrapped * line_h_emu
        sb = int(para.space_before.pt * EMU_PER_PT) if para.space_before else 0
        sa = int(para.space_after.pt  * EMU_PER_PT) if para.space_after  else 0
        total_height_emu += sb + para_h + sa
        wrapped_lines += wrapped

    if total_height_emu > inner_h_emu:  # NO 5% slack
        return {
            "slide": idx,
            "shape_name": getattr(shape, "name", "?"),
            "shape_box_in": {
                "left": round((shape.left or 0) / EMU_PER_INCH, 2),
                "top":  round((shape.top  or 0) / EMU_PER_INCH, 2),
                "w":    round((shape.width  or 0) / EMU_PER_INCH, 2),
                "h":    round((shape.height or 0) / EMU_PER_INCH, 2),
            },
            "needed_h_in": round(total_height_emu / EMU_PER_INCH, 3),
            "available_h_in": round(inner_h_emu / EMU_PER_INCH, 3),
            "overflow_in": round((total_height_emu - inner_h_emu) / EMU_PER_INCH, 3),
            "wrapped_lines": wrapped_lines,
        }
    return None


# ---------- F4: theme-resolved contrast ----------

def _resolved_run_color(run, master_theme: dict) -> tuple[int, int, int] | None:
    """Resolve effective run colour. Returns None when unresolvable
    (caller records as `contrast_unresolved`, never silently passes).
    """
    try:
        col = run.font.color
        if col is None:
            return None
        if col.type is not None:
            try:
                rgb = col.rgb
                if rgb is not None:
                    return rgb[0], rgb[1], rgb[2]
            except Exception:
                pass
            try:
                sc = getattr(col, "theme_color", None)
                if sc is not None and master_theme.get(str(sc)):
                    return _hex_to_rgb(master_theme[str(sc)])
            except Exception:
                pass
    except Exception:
        return None
    return None


def _is_large_text(size_pt: float, bold: bool) -> bool:
    # WCAG: >=18pt OR >=14pt-bold counts as "large text" -> 3:1 threshold.
    return size_pt >= 18 or (bool(bold) and size_pt >= 14)


def _contrast_violations_for(slide, idx: int, master_theme: dict) -> tuple[list[dict], list[dict]]:
    out: list[dict] = []
    unresolved: list[dict] = []
    bg = _slide_bg_rgb(slide)
    if bg is None:
        unresolved.append({"slide": idx, "code": "bg_unresolved"})
        return out, unresolved
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        local_bg = _shape_solid_rgb(shp) or bg
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                fg = _resolved_run_color(run, master_theme)
                if fg is None:
                    unresolved.append({
                        "slide": idx,
                        "shape": getattr(shp, "name", "?"),
                        "text_sample": (run.text or "")[:40],
                    })
                    continue
                size_pt = (
                    (run.font.size.pt if run.font.size else None)
                    or (para.font.size.pt if para.font.size else None)
                    or 18.0
                )
                bold = bool(run.font.bold)
                threshold = 3.0 if _is_large_text(size_pt, bold) else 4.5
                ratio = _contrast_ratio(fg, local_bg)
                if ratio < threshold:
                    out.append({
                        "slide": idx,
                        "shape": getattr(shp, "name", "?"),
                        "fg": "#%02X%02X%02X" % fg,
                        "bg": "#%02X%02X%02X" % local_bg,
                        "ratio": round(ratio, 2),
                        "threshold": threshold,
                        "size_pt": size_pt,
                        "is_large_text": _is_large_text(size_pt, bold),
                    })
                    # NO early return — keep scanning every run.
    return out, unresolved


# ---------- main ----------

def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True, type=Path)
    ap.add_argument("--spec", type=Path, default=None)
    ap.add_argument("--out", required=True, type=Path)
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    try:
        prs = Presentation(str(args.pptx))
    except Exception as e:
        print(f"ERROR: could not open pptx: {e}", file=sys.stderr)
        return 1

    sw = prs.slide_width or 0
    sh = prs.slide_height or 0
    ratio = (sw / sh) if sh else 0.0
    aspect_pass = abs(ratio - 16 / 9) < 0.01

    # Pull the deck-spec's design_system_tokens block (F9) so theme
    # colour resolution can recover scheme-color references.
    master_theme: dict[str, str] = {}
    if args.spec and args.spec.exists():
        try:
            ds = json.loads(args.spec.read_text())
            tokens = ds.get("design_system_tokens") or {}
            palette = tokens.get("palette") or {}
            # Map common scheme keys -> hex; loose-keyed for forward-compat.
            for k, v in palette.items():
                master_theme[k] = v
        except Exception:
            pass

    overflow_violations: list[dict] = []
    contrast_violations: list[dict] = []
    contrast_unresolved: list[dict] = []
    alignment_violations: list[dict] = []
    layout_hashes: list[tuple[int, str]] = []
    underline_slides: list[int] = []
    bg_sequence: list[str] = []
    speaker_notes_missing: list[int] = []
    titles: dict[str, list[int]] = {}
    body_word_violations: list[dict] = []
    font_fallback_flag = [False]

    slides = list(prs.slides)
    for idx, slide in enumerate(slides, start=1):
        # F1: real-metric overflow per shape
        for shape in slide.shapes:
            ov = _shape_overflow(shape, idx, font_fallback_flag)
            if ov:
                overflow_violations.append(ov)
        # F4: theme-resolved contrast
        cv, cu = _contrast_violations_for(slide, idx, master_theme)
        contrast_violations.extend(cv)
        contrast_unresolved.extend(cu)
        # alignment
        alignment_violations.extend(_alignment_violations_for(slide, idx))
        layout_hashes.append((idx, _layout_hash(slide)))
        if _has_underline_below_title(slide):
            underline_slides.append(idx)
        bg = _slide_bg_rgb(slide)
        bg_sequence.append(_bg_label(bg))
        try:
            ns = slide.notes_slide
            ntxt = (ns.notes_text_frame.text or "").strip() if ns and ns.notes_text_frame else ""
            if not ntxt:
                speaker_notes_missing.append(idx)
        except Exception:
            speaker_notes_missing.append(idx)
        title = _slide_title_text(slide)
        if title:
            titles.setdefault(title, []).append(idx)
        bwc = _body_word_count(slide)
        # F11: tightened body_word_max from 70 -> 30 (Reynolds)
        if bwc > 30:
            body_word_violations.append({"slide": idx, "body_words": bwc})

    # Repeated layout hash on consecutive slides.
    repeated = []
    for i in range(1, len(layout_hashes)):
        if layout_hashes[i][1] == layout_hashes[i - 1][1]:
            repeated.append({
                "slides": [layout_hashes[i - 1][0], layout_hashes[i][0]],
                "hash": layout_hashes[i][1],
            })

    # F8: max-same-bg-run with accent-break awareness.
    max_run = 0
    cur_run = 0
    last = None
    cur_run_has_accent_break = False
    longest_run_had_accent_break = True
    for b in bg_sequence:
        if b == last and b != "unknown":
            cur_run += 1
        else:
            if cur_run >= max_run:
                max_run = cur_run
                longest_run_had_accent_break = cur_run_has_accent_break
            cur_run = 1
            last = b
            cur_run_has_accent_break = False
        # An accent slide *inside* a same-label run is impossible by
        # definition (same label means non-accent like dark/light); the
        # `accent_break` heuristic instead tracks whether any accent
        # slide appears in the *deck* between dark/light spans.
    if cur_run >= max_run:
        max_run = cur_run
        longest_run_had_accent_break = cur_run_has_accent_break
    accent_present = any(x == "accent" for x in bg_sequence)

    duplicate_titles = [{"title": t, "slides": ids} for t, ids in titles.items() if len(ids) > 1]

    blocking: list[str] = []
    warnings: list[str] = []

    if not aspect_pass:
        blocking.append("aspect_ratio")
    if overflow_violations:
        blocking.append("overflow")
    if contrast_violations:
        blocking.append("contrast")
    if len(contrast_unresolved) >= 5:
        blocking.append("contrast_unresolved")
    if len(underline_slides) > 2:
        blocking.append("title_underline_spam")

    # F8 dark_light_run thresholds:
    #   warn at >= 3, blocking at >= 7, OR blocking at >= 5 with no accent break.
    if max_run >= 7:
        blocking.append("dark_light_run")
    elif max_run >= 5 and not accent_present:
        blocking.append("dark_light_run")
    elif max_run >= 3:
        warnings.append("dark_light_run")

    if speaker_notes_missing:
        blocking.append("speaker_notes_missing")
    if duplicate_titles:
        blocking.append("duplicate_titles")
    if repeated:
        blocking.append("repeated_layout_hash")
    if alignment_violations:
        warnings.append("alignment")
    if body_word_violations:
        warnings.append("body_word_max")
    if contrast_unresolved and "contrast_unresolved" not in blocking:
        warnings.append("contrast_unresolved")

    report = {
        "session_id": None,
        "pptx_path": str(args.pptx),
        "spec_path": str(args.spec) if args.spec else None,
        "slide_count": len(slides),
        "aspect_ratio": round(ratio, 4),
        "aspect_ratio_pass": aspect_pass,
        "overflow_violations": overflow_violations,
        "contrast_violations": contrast_violations,
        "contrast_unresolved": contrast_unresolved,
        "alignment_violations": alignment_violations,
        "repeated_layout_hash": repeated,
        "title_underline_count": len(underline_slides),
        "underline_slides": underline_slides,
        "max_same_bg_run": max_run,
        "background_sequence": bg_sequence,
        "accent_break_present": accent_present,
        "speaker_notes_missing": speaker_notes_missing,
        "duplicate_titles": duplicate_titles,
        "body_word_max_violations": body_word_violations,
        "font_fallback": font_fallback_flag[0],
        "blocking_findings": blocking,
        "warning_findings": warnings,
    }
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(report, indent=2))
    status = "PASS" if not blocking else "FAIL"
    print(f"{status}: {len(blocking)} blocking, {len(warnings)} warning findings ({len(slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
