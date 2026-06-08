#!/usr/bin/env python3
"""generate_deck.py — Token-driven deck builder consumed by @deck-builder.

Pipeline (drives off `deck-spec.json`):
  1. Load `deck-spec.json` (validate against schema is the agent's job).
  2. Resolve the design system: prefer `deck-spec.design_system_tokens`
     (F9 — token-driven constants); fall back to a sensible default
     palette when the spec doesn't supply tokens.
  3. Pre-render every `slide.visual_assets[]` entry by subprocess'ing
     the matching `render-visual` script (chart / composite / diagram).
     Failures of render-visual diagrams (graceful-degrade per OQ2)
     produce `<out>.skipped.json` sentinels — we treat these as
     "asset unavailable; build slide without it" and surface the
     skip in the build log.
  4. Dispatch each slide on `slide.style`:
       - "simple" (default — C5 backwards-compat) → simple builder
       - "styled" → styled builder selected by `slide.style_recipe`
     Unknown / inconsistent style or recipe → reject with
     `bad_spec.<reason>` and exit 2.
  5. Write `output.pptx`. Print SUCCESS / ERROR.

Usage:
  python generate_deck.py --spec deck-spec.json --out output.pptx
                          [--template path/to/template.pptx]
                          [--build-log build-log.txt]

Backwards-compat (C5):
  - Slides without a `style` key are treated as `style: "simple"`.
  - Decks without `design_system_tokens` use the fallback palette.
  - Existing simple types (title, key-message, content, ...)
    continue to work unchanged.
"""

from __future__ import annotations

import argparse
import json
import os
import subprocess
import sys
from pathlib import Path

# --- Dependency check ---
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE


# --- 16:9 EMU constants (architecture §4.1) ---
SLIDE_W_EMU = 12192000   # 13.333"
SLIDE_H_EMU = 6858000    # 7.5"
SLIDE_HALF_EMU = 6096000

# --- Default fallback palette (C5: when deck-spec lacks tokens) ---
FALLBACK_TOKENS: dict = {
    "name": "fallback",
    "palette": {
        "background_dark":   "#0F1B2D",
        "background_light":  "#F4F5F7",
        "background_accent": "#3B82F6",
        "primary_accent":    "#3B82F6",
        "secondary_accent":  "#06B6D4",
        "highlight":         "#F59E0B",
        "surface_elevated":  "#FFFFFF",
        "text_on_dark":      "#FFFFFF",
        "text_on_light":     "#1A1A2E",
        "text_secondary":    "#6B7080",
        "text_secondary_on_dark":  "#6B7080",
        "text_secondary_on_light": "#6B7080",
    },
    "typography": {
        "font_title": "Inter",
        "font_body":  "Inter",
        "title_weight": 700,
        "body_weight":  400,
        "size_hero":     54,
        "size_section":  48,
        "size_title":    40,
        "size_subtitle": 24,
        "size_body":     22,
        "size_caption":  14,
        "size_metric_xxl": 200,
        "size_quote_glyph": 240,
    },
    "chart_palette": {
        "focal": "#3B82F6",
        "muted": "#9AA1AC",
        "grid":  "#D8DBE0",
        "ramp": ["#3B82F6", "#06B6D4", "#8B5CF6", "#F59E0B", "#10B981", "#EF4444"],
    },
}

# --- Canonical recipe sets (architecture §5) ---
STYLED_RECIPES = {
    "hero_full_bleed", "accent_block_left", "metric_xxl",
    "quote_pullout", "split_image_right", "tinted_panel_right",
    "progress_dots", "chart_callout",
    # Session 2026-05-04-c8d3b2a1 — archetype builders
    "risk_heatmap", "priority_matrix", "waterfall",
    "flywheel", "funnel", "decision_options",
    "appendix_dense",
    # Session 2026-06-08-c5d9e1a7 (C6) — high-impact editorial archetypes
    "stat_grid_3up", "pull_quote_portrait", "full_bleed_caption",
    "editorial_2col_6040", "timeline_horizontal", "agenda_toc",
    "closing_cta",
    # `footer_source` is a partial applied via slide.footer; not a primary recipe.
}
SIMPLE_TYPES = {
    "title", "key-message", "content", "metric-spotlight",
    "comparison-columns", "quote", "data-callout",
    "section-divider", "visual-hero", "question", "cta-steps",
}


# ============================================================
# Token resolution helpers
# ============================================================

def _hex(s: str) -> RGBColor:
    s = (s or "#000000").lstrip("#")
    if len(s) != 6:
        s = "000000"
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _rgb_tuple(c: RGBColor) -> tuple[int, int, int]:
    return (c[0], c[1], c[2])


def _mix(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
    """Linear blend a→b by t in [0,1]. t=0 → a, t=1 → b."""
    t = max(0.0, min(1.0, t))
    return RGBColor(
        int(round(a[0] + (b[0] - a[0]) * t)),
        int(round(a[1] + (b[1] - a[1]) * t)),
        int(round(a[2] + (b[2] - a[2]) * t)),
    )


_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_BLACK = RGBColor(0x00, 0x00, 0x00)


def _tint(c: RGBColor, amount: float) -> RGBColor:
    """Lighten toward white by `amount` (Refactoring-UI tint ladder)."""
    return _mix(c, _WHITE, amount)


def _shade(c: RGBColor, amount: float) -> RGBColor:
    """Darken toward black by `amount` (Refactoring-UI shade ladder)."""
    return _mix(c, _BLACK, amount)


# ------------------------------------------------------------
# WCAG contrast auto-safety (mirrors pptx-structural-asserts
# check_pptx: _rgb_to_relative_luminance / _contrast_ratio /
# _is_large_text / _resolved_run_color / _shape_solid_rgb /
# _slide_bg_rgb). Applied as a final pass before save so accent
# text marks (delta chips, node/date labels) that would otherwise
# fail AA on a light surface are nudged darker (or lighter on a dark
# surface) just enough to resolve — keeping hue, improving legibility,
# and letting the visual-QA loop reach pass with *genuine* contrast
# rather than masked/false-blocked contrast.
# ------------------------------------------------------------

def _rel_lum(rgb: tuple[int, int, int]) -> float:
    def chan(c: int) -> float:
        x = c / 255.0
        return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    l1, l2 = _rel_lum(fg), _rel_lum(bg)
    a, b = max(l1, l2), min(l1, l2)
    return (a + 0.05) / (b + 0.05)


def _is_large(size_pt: float, bold: bool) -> bool:
    return size_pt >= 18 or (bool(bold) and size_pt >= 14)


def _shape_fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        f = shape.fill
        if f.type is None:
            return None
        if hasattr(f, "fore_color") and getattr(f.fore_color, "type", None) is not None:
            rgb = f.fore_color.rgb
            if rgb is not None:
                return rgb[0], rgb[1], rgb[2]
    except Exception:
        return None
    return None


def _slide_bg_tuple(slide) -> tuple[int, int, int] | None:
    try:
        bg = slide.background
        if bg is not None and bg.fill is not None and getattr(bg.fill, "type", None) is not None:
            rgb = bg.fill.fore_color.rgb
            if rgb is not None:
                return rgb[0], rgb[1], rgb[2]
    except Exception:
        pass
    best, best_area = None, 0
    for shp in slide.shapes:
        try:
            area = int(shp.width or 0) * int(shp.height or 0)
            rgb = _shape_fill_rgb(shp)
            if rgb and area > best_area:
                best, best_area = rgb, area
        except Exception:
            continue
    return best


def _aa_safe(fg: tuple[int, int, int], bg: tuple[int, int, int],
             threshold: float) -> tuple[int, int, int] | None:
    """Return an AA-safe variant of `fg` against `bg` (>= threshold), or
    None if `fg` already passes. Preserves hue by mixing toward the pole
    (black/white) that raises contrast, stepping until the bar is met."""
    if _contrast(fg, bg) >= threshold:
        return None
    target = _BLACK if _rel_lum(bg) > _rel_lum(fg) else _WHITE
    fc = RGBColor(*fg)
    for step in range(1, 21):
        cand = _mix(fc, target, step / 20.0)
        if _contrast((cand[0], cand[1], cand[2]), bg) >= threshold:
            return (cand[0], cand[1], cand[2])
    return (target[0], target[1], target[2])


def _enforce_text_contrast(prs, log=None) -> int:
    """Final-pass WCAG guard: darken/lighten only the text runs that fail
    AA against their effective background. Returns the count of runs fixed."""
    fixed = 0
    for slide in prs.slides:
        bg = _slide_bg_tuple(slide)
        if bg is None:
            continue
        for shp in slide.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            local_bg = _shape_fill_rgb(shp) or bg
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        col = run.font.color
                        if col is None or col.type is None:
                            continue
                        rgb = col.rgb
                        if rgb is None:
                            continue
                    except Exception:
                        continue
                    fg = (rgb[0], rgb[1], rgb[2])
                    size_pt = (
                        (run.font.size.pt if run.font.size else None)
                        or (para.font.size.pt if para.font.size else None)
                        or 18.0
                    )
                    bold = bool(run.font.bold)
                    threshold = 3.0 if _is_large(size_pt, bold) else 4.5
                    safe = _aa_safe(fg, local_bg, threshold)
                    if safe is not None:
                        run.font.color.rgb = RGBColor(*safe)
                        fixed += 1
    if log is not None and fixed:
        log.write(f"[contrast] auto-corrected {fixed} sub-AA text run(s)\n")
    return fixed


class Tokens:
    """Token-driven palette + typography accessor (F9)."""

    def __init__(self, raw: dict):
        self.raw = raw
        p = raw.get("palette", {})
        t = raw.get("typography", {})
        # Palette colours (RGBColor)
        self.bg_dark    = _hex(p.get("background_dark",   "#0F1B2D"))
        self.bg_light   = _hex(p.get("background_light",  "#F4F5F7"))
        self.bg_accent  = _hex(p.get("background_accent", "#3B82F6"))
        self.primary    = _hex(p.get("primary_accent",    "#3B82F6"))
        self.secondary  = _hex(p.get("secondary_accent",  "#06B6D4"))
        self.highlight  = _hex(p.get("highlight",         "#F59E0B"))
        self.surface    = _hex(p.get("surface_elevated",  "#FFFFFF"))
        self.text_dark  = _hex(p.get("text_on_dark",      "#FFFFFF"))
        self.text_light = _hex(p.get("text_on_light",     "#1A1A2E"))
        # Secondary/muted text. Some systems (e.g. technical-slate) ship
        # surface-scoped overrides because no single mid-gray clears AA
        # against both their dark and light backgrounds. Prefer the
        # override when present, else fall back to the legacy bare key.
        _ts_default = p.get("text_secondary", "#6B7080")
        self.text_2_dark  = _hex(p.get("text_secondary_on_dark",  _ts_default))
        self.text_2_light = _hex(p.get("text_secondary_on_light", _ts_default))
        # Backward-compat alias (dark surface dominates most decks).
        self.text_2     = self.text_2_dark
        # Typography
        self.font_title = t.get("font_title", "Inter")
        self.font_body  = t.get("font_body",  "Inter")
        self.sz_hero     = Pt(t.get("size_hero", 54))
        self.sz_section  = Pt(t.get("size_section", 48))
        self.sz_title    = Pt(t.get("size_title", 40))
        self.sz_subtitle = Pt(t.get("size_subtitle", 24))
        self.sz_body     = Pt(t.get("size_body", 22))
        self.sz_caption  = Pt(t.get("size_caption", 14))
        self.sz_metric_xxl  = Pt(t.get("size_metric_xxl", 200))
        self.sz_quote_glyph = Pt(t.get("size_quote_glyph", 240))

        # --- C3: font-weight tokens (300/400/600/700/800). python-pptx
        # cannot set arbitrary numeric weights directly, so helpers map
        # weight >= 600 to bold. Stored for builders/_set_tracking use.
        self.title_weight = int(t.get("title_weight", 700))
        self.body_weight  = int(t.get("body_weight", 400))
        self.eyebrow_weight = int(t.get("eyebrow_weight", 600))

        # --- C3: tint / shade ladders. Builders layer tone (cards,
        # insets, hairlines, scrims) from these. Explicit ladders in
        # `palette.tints` win; otherwise derived from base colours so
        # legacy specs + the fallback palette always have tones.
        tints = p.get("tints") or {}
        self._tints = {
            role: {str(k): _hex(v) for k, v in ladder.items()}
            for role, ladder in tints.items() if isinstance(ladder, dict)
        }
        # Craft-surface derived tones (overridable via palette keys).
        # `hairline` — 1px rule colour: a faint mix toward the text.
        self.hairline = _hex(p["hairline"]) if p.get("hairline") else \
            _mix(self.bg_light, self.text_light, 0.14)
        self.hairline_on_dark = _hex(p["hairline_on_dark"]) if p.get("hairline_on_dark") else \
            _mix(self.bg_dark, self.text_dark, 0.18)
        # `card_fill` — elevated card surface on a light slide.
        self.card_fill = self.surface if p.get("surface_elevated") else \
            _tint(self.bg_light, 0.45)
        self.card_fill_on_dark = _hex(p["surface_on_dark"]) if p.get("surface_on_dark") else \
            _tint(self.bg_dark, 0.08)
        # `scrim` — overlay colour for full-bleed caption contrast.
        self.scrim = _hex(p["scrim"]) if p.get("scrim") else self.bg_dark

        # --- C3: categorical chart palette. Multi-series charts cycle
        # `ramp`; the focal series uses `chart_focal`; muted marks use
        # `chart_muted`; gridlines/axes use `chart_grid`.
        cp = raw.get("chart_palette") or {}
        self.chart_focal = _hex(cp.get("focal")) if cp.get("focal") else self.secondary
        self.chart_muted = _hex(cp.get("muted")) if cp.get("muted") else self.text_2_dark
        self.chart_grid  = _hex(cp.get("grid")) if cp.get("grid") else self.text_2_light
        ramp = cp.get("ramp") or []
        self.chart_ramp = [_hex(c) for c in ramp] if ramp else [
            self.primary, self.secondary, self.highlight,
            _tint(self.primary, 0.35), _shade(self.secondary, 0.25),
            _tint(self.highlight, 0.3),
        ]

    def tint(self, role: str, step: str, default: RGBColor = None) -> RGBColor:
        """Look up an explicit tint-ladder colour (e.g. tint('primary','100'))."""
        ladder = self._tints.get(role) or {}
        if step in ladder:
            return ladder[step]
        return default if default is not None else self.primary


# ============================================================
# Low-level shape helpers
# ============================================================

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_notes(slide, notes: str) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_rect(slide, x, y, w, h, color: RGBColor):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _add_textbox(slide, x, y, w, h, text, *,
                 font_name="Calibri", font_size=Pt(22),
                 color: RGBColor = None, bold=False, align=None,
                 anchor=None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = font_size
        p.font.bold = bold
        if color is not None:
            # Paragraph-level set kept for inherited defaults, but the
            # contrast checker (check_pptx._resolved_run_color) reads colour
            # at RUN level. Set colour on every run so real contrast resolves
            # instead of bucketing as `contrast_unresolved`.
            p.font.color.rgb = color
            for r in p.runs:
                r.font.name = font_name
                r.font.size = font_size
                r.font.bold = bold
                r.font.color.rgb = color
        if align is not None:
            p.alignment = align
    return tx


def _add_picture(slide, path, x, y, w, h):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w, height=h)


def _new_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ============================================================
# C1 — Rendering-craft helpers (the leverage move)
# ============================================================
#
# Four reusable primitives that turn flat solid()-fill blocks into
# layered, "premium"-feeling compositions: tonal cards, hairline rules,
# gradient scrims, and letter-tracking on display type. Used across the
# styled + archetype builders.

from pptx.oxml.ns import qn  # noqa: E402


def _set_alpha(fill_or_color_elem, pct: int) -> None:
    """Inject an <a:alpha val=".."/> into a solidFill's srgbClr."""
    try:
        srgb = fill_or_color_elem.find(qn("a:srgbClr"))
        if srgb is None:
            return
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
        srgb.append(alpha)
    except Exception:
        pass


def _add_soft_shadow(shape, blur_emu: int = 90000, dist_emu: int = 38100,
                     alpha_pct: int = 72) -> None:
    """Attach a subtle outer drop shadow to a shape (manual, not inherited).

    `alpha_pct` is the shadow's TRANSPARENCY (higher = fainter)."""
    try:
        spPr = shape._element.spPr
        # Remove any inherited effect list first.
        for tag in ("a:effectLst",):
            existing = spPr.find(qn(tag))
            if existing is not None:
                spPr.remove(existing)
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        shdw = effectLst.makeelement(qn("a:outerShdw"), {
            "blurRad": str(int(blur_emu)),
            "dist": str(int(dist_emu)),
            "dir": "5400000",   # straight down
            "rotWithShape": "0",
        })
        clr = shdw.makeelement(qn("a:srgbClr"), {"val": "1A1A1A"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha_pct) * 1000))})
        clr.append(alpha)
        shdw.append(clr)
        effectLst.append(shdw)
        spPr.append(effectLst)
    except Exception:
        pass


def _add_card(slide, x, y, w, h, *, fill: RGBColor, line: RGBColor = None,
              line_w_emu: int = 9525, radius: float = 0.06,
              shadow: bool = True):
    """Rounded-rect 'card': tonal fill + hairline border + soft shadow.

    `radius` is the rounded-corner adjustment (0..0.5 of the short side).
    Returns the shape so callers can place content over it."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line is not None:
        card.line.color.rgb = line
        card.line.width = Emu(int(line_w_emu))
    else:
        card.line.fill.background()
    # Corner radius via the shape adjustment handle.
    try:
        card.adjustments[0] = radius
    except Exception:
        pass
    if shadow:
        _add_soft_shadow(card)
    return card


def _add_hairline(slide, x, y, w, color: RGBColor, *, weight_emu: int = 9525,
                  vertical: bool = False, length=None):
    """Thin 1px rule (horizontal by default). The single most effective
    'editorial' primitive — used under kickers, between rows, as folio
    separators. `weight_emu` 9525 ≈ 0.75pt."""
    if vertical:
        h = length if length is not None else w
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(int(weight_emu)), h)
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(weight_emu)))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def _add_scrim(slide, x, y, w, h, color: RGBColor, *,
               top_alpha: int = 0, bottom_alpha: int = 78,
               direction: int = 5400000):
    """Vertical gradient scrim (transparent → opaque `color`) guaranteeing
    caption/title contrast over full-bleed imagery. `*_alpha` are OPACITY
    percentages (0 = fully transparent, 100 = solid) at each end."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.line.fill.background()
    rect.shadow.inherit = False
    try:
        spPr = rect._element.spPr
        # Drop the solidFill python-pptx created by default.
        for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            ex = spPr.find(qn(tag))
            if ex is not None:
                spPr.remove(ex)
        hexv = "%02X%02X%02X" % (color[0], color[1], color[2])
        grad = spPr.makeelement(qn("a:gradFill"), {})
        gsLst = grad.makeelement(qn("a:gsLst"), {})
        for pos, op in ((0, top_alpha), (100000, bottom_alpha)):
            gs = gsLst.makeelement(qn("a:gs"), {"pos": str(pos)})
            clr = gs.makeelement(qn("a:srgbClr"), {"val": hexv})
            a = clr.makeelement(qn("a:alpha"), {"val": str(int(op * 1000))})
            clr.append(a)
            gs.append(clr)
            gsLst.append(gs)
        grad.append(gsLst)
        lin = grad.makeelement(qn("a:lin"), {"ang": str(direction), "scaled": "1"})
        grad.append(lin)
        # Insert gradFill in the correct schema position (after line? before).
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            ln.addprevious(grad)
        else:
            spPr.append(grad)
    except Exception:
        # Fallback: plain semi-opaque solid.
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
    return rect


def _set_tracking(textbox, spc: int) -> None:
    """Apply letter-spacing (tracking) to every run in a textbox.
    `spc` is in 1/100 pt (e.g. 200 = +2pt). Negative tightens display
    type; positive (+120..+300) opens tracked uppercase kickers."""
    try:
        txBody = textbox.text_frame._txBody
        for p in txBody.findall(qn("a:p")):
            for r in p.findall(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is None:
                    rPr = r.makeelement(qn("a:rPr"), {})
                    r.insert(0, rPr)
                rPr.set("spc", str(int(spc)))
    except Exception:
        pass


def _add_arrowhead(connector, tail: bool = True, head: bool = False) -> None:
    """Add triangle line-end(s) to a connector/line shape via DrawingML."""
    try:
        ln = connector.line._get_or_add_ln()
        if tail:
            te = ln.makeelement(qn("a:tailEnd"),
                                {"type": "triangle", "w": "med", "len": "med"})
            ln.append(te)
        if head:
            he = ln.makeelement(qn("a:headEnd"),
                                {"type": "triangle", "w": "med", "len": "med"})
            ln.append(he)
    except Exception:
        pass


def _add_kicker(slide, x, y, w, text, t, color, *,
                rule: bool = True, rule_color: RGBColor = None):
    """Tracked uppercase eyebrow + optional hairline rule beneath it.
    The editorial signature used across premium builders."""
    if not text:
        return None
    # Kickers are bold uppercase display eyebrows. Floor the size at the
    # WCAG large-text boundary (14pt bold) so an accent-coloured eyebrow is
    # evaluated at the 3:1 threshold, not 4.5:1 — a 13pt caption token (e.g.
    # ink-editorial) otherwise pushes a legible accent eyebrow into a false
    # contrast block. 14pt bold tracked is the canonical editorial eyebrow.
    kicker_size = t.sz_caption if t.sz_caption >= Pt(14) else Pt(14)
    tx = _add_textbox(slide, x, y, w, Inches(0.4),
                      str(text).upper(),
                      font_name=t.font_body, font_size=kicker_size,
                      color=color, bold=True)
    _set_tracking(tx, 220)
    if rule:
        _add_hairline(slide, x, y + Inches(0.42), Inches(0.9),
                      rule_color if rule_color is not None else color,
                      weight_emu=19050)
    return tx


# ============================================================
# Pre-render visual_assets[]
# ============================================================

def _prerender_assets(spec_path: Path, deck_spec: dict, log) -> None:
    """For each slide.visual_assets[], subprocess the appropriate
    render-visual script. Mutate the slide dict in place — replace each
    asset entry's `path` field with the produced PNG path.

    Each asset entry shape (additive):
      {"kind": "chart" | "composite" | "diagram",
       "recipe": "bar_with_callouts" | "gradient_pattern" | ...,
       "spec": {...},                  -- recipe-specific
       "out": "path/to/asset.png",     -- where to write
       "size": "1920x1080",            -- optional
       "svg": false}                   -- optional
    """
    here = Path(__file__).resolve().parent.parent.parent.parent  # .github/skills/
    rv_dir = here / "render-visual" / "scripts"
    script_for = {
        "chart":     rv_dir / "render_chart.py",
        "composite": rv_dir / "render_composite.py",
        "diagram":   rv_dir / "render_diagram.py",
    }

    tokens_path = spec_path.parent / "_tokens.json"
    tokens_path.write_text(json.dumps(
        deck_spec.get("design_system_tokens") or FALLBACK_TOKENS, indent=2))

    for slide in deck_spec.get("slides", []):
        for asset in slide.get("visual_assets") or []:
            kind = asset.get("kind")
            recipe = asset.get("recipe")
            out = asset.get("out")
            if kind not in script_for or not recipe or not out:
                log.write(f"[asset] skip (bad kind/recipe/out): {asset}\n")
                continue
            script = script_for[kind]
            if not script.exists():
                log.write(f"[asset] script missing: {script}\n")
                continue
            spec_file = spec_path.parent / f"_asset_{slide.get('index','?')}_{recipe}.json"
            spec_file.write_text(json.dumps(asset.get("spec") or {}, indent=2))
            cmd = [
                sys.executable, str(script),
                "--kind", recipe,
                "--spec", str(spec_file),
                "--tokens", str(tokens_path),
                "--out", str(out),
                "--size", asset.get("size", "1920x1080"),
            ]
            if asset.get("svg"):
                cmd.append("--svg")
            try:
                r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                log.write(f"[asset] {recipe} rc={r.returncode}\n")
                if r.stdout:
                    log.write(r.stdout + "\n")
                if r.stderr:
                    log.write(r.stderr + "\n")
                # If the script produced a .skipped.json sentinel, mark it.
                sentinel = Path(str(out) + ".skipped.json")
                if sentinel.exists():
                    asset["skipped"] = True
                    log.write(f"[asset] {recipe} skipped (graceful-degrade)\n")
                else:
                    asset["path"] = out
            except Exception as e:
                log.write(f"[asset] {recipe} exception: {e}\n")
                asset["skipped"] = True


# ============================================================
# SIMPLE BUILDERS  (per slide.type when style == "simple")
# ============================================================

def _simple_title(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_rect(s, 0, 0, SLIDE_W_EMU, Emu(54864), t.primary)  # 0.06" top bar
    # Editorial eyebrow + hairline above the hero (asymmetric, anchored left).
    _add_kicker(s, Inches(1.2), Inches(1.5), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "",
                t, t.text_2_dark, rule_color=t.primary)
    title_tx = _add_textbox(s, Inches(1.2), Inches(2.3), Inches(10.9), Inches(1.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_hero,
                 color=t.text_dark)
    _set_tracking(title_tx, -20)  # tighten display type
    if slide.get("subtitle"):
        _add_textbox(s, Inches(1.2), Inches(4.4), Inches(10.9), Inches(0.8),
                     slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.secondary)
    # Footer rule + meta (author / date / confidential) when supplied.
    meta = slide.get("meta") or slide.get("footer_meta")
    if meta:
        _add_hairline(s, Inches(1.2), Inches(6.5), Inches(10.9), t.hairline_on_dark)
        _add_textbox(s, Inches(1.2), Inches(6.65), Inches(10.9), Inches(0.5),
                     str(meta), font_name=t.font_body, font_size=t.sz_caption,
                     color=t.text_2_dark)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_section_divider(prs, slide, t: Tokens):
    """C5 backwards-compat path. NOTE: per OQ4, the strategist defaults
    section dividers to `style: styled / hero_full_bleed`; this simple
    path runs only when the spec explicitly opts back to simple."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_rect(s, Inches(5.5), Inches(2.5), Inches(2.3), Emu(32004), t.secondary)
    _add_textbox(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_section,
                 color=t.text_dark, align=PP_ALIGN.CENTER)
    if slide.get("subtitle"):
        _add_textbox(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
                     slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_2_dark, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_key_message(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(1.2), Inches(2.5), Inches(10.9), Inches(2.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_hero,
                 color=t.text_light, bold=False)
    if slide.get("subtitle"):
        _add_textbox(s, Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.6),
                     slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_2_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_content(prs, slide, t: Tokens):
    """C2 — editorial lead + supporting points, NOT a bullet dump.

    Layout: tracked kicker, action title, an optional bold *lead*
    standfirst, then <=3 supporting points rendered as tonal cards in a
    right-weighted 2-col grid (or a single column for <=2 points). The
    flat full-height `"•  {b}"` textbox is gone — that was antipattern-3
    incarnate and it was the DEFAULT content layout."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_kicker(s, Inches(0.75), Inches(0.55), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "",
                t, t.text_2_light, rule_color=t.primary)
    _add_textbox(s, Inches(0.75), Inches(1.05), Inches(11.83), Inches(1.0),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    bullets = slide.get("bullets") or []
    lead = slide.get("lead") or slide.get("standfirst")
    # Back-compat: if no explicit lead and there are bullets, promote the
    # first bullet to the lead so legacy specs still read editorially.
    if not lead and bullets:
        lead, bullets = bullets[0], bullets[1:]
    points = list(bullets)[:3]  # hard cap at 3 supporting points

    y = Inches(2.15)
    if lead:
        lead_tx = _add_textbox(s, Inches(0.75), y, Inches(11.83), Inches(1.2),
                     str(lead),
                     font_name=t.font_title, font_size=Pt(26),
                     color=t.text_light, bold=False)
        y = Inches(3.45)

    if not points:
        _set_notes(s, slide.get("notes", ""))
        return s

    # Tonal cards: 2-up for 3-4 points, single-col for <=2.
    cols = 1 if len(points) <= 2 else 2
    gutter = Inches(0.3)
    body_left = Inches(0.75)
    body_w = Inches(11.83)
    card_w = (body_w - gutter * (cols - 1)) / cols
    rows = (len(points) + cols - 1) // cols
    avail_h = Inches(6.6) - y
    card_h = (avail_h - gutter * (rows - 1)) / rows
    for i, pt in enumerate(points):
        ci = i % cols
        ri = i // cols
        cx = body_left + (card_w + gutter) * ci
        cy = y + (card_h + gutter) * ri
        _add_card(s, cx, cy, card_w, card_h, fill=t.card_fill,
                  line=t.hairline, shadow=True)
        # Accent index numeral + point text.
        _add_textbox(s, cx + Inches(0.3), cy + Inches(0.2),
                     Inches(0.8), Inches(0.6), f"{i + 1:02d}",
                     font_name=t.font_title, font_size=Pt(22),
                     color=t.primary, bold=True)
        _add_textbox(s, cx + Inches(0.3), cy + Inches(0.85),
                     card_w - Inches(0.6), card_h - Inches(1.0),
                     str(pt), font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light, anchor=MSO_ANCHOR.TOP)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_metric_spotlight(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    metric = str(slide.get("metric", ""))
    label  = slide.get("metric_label", "")
    _add_textbox(s, Inches(0.75), Inches(1.8), Inches(11.83), Inches(3.5),
                 metric,
                 font_name=t.font_title, font_size=Pt(150),
                 color=t.primary, align=PP_ALIGN.CENTER, bold=True)
    if label:
        _add_textbox(s, Inches(0.75), Inches(5.4), Inches(11.83), Inches(0.7),
                     label,
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_2_light, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_comparison_columns(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    cols = slide.get("columns") or []
    col_w = Inches(11.0 / max(1, len(cols)))
    for i, col in enumerate(cols):
        x = Inches(0.75) + col_w * i
        _add_rect(s, x, Inches(1.8), Emu(32004), Inches(4.5), t.primary)
        _add_textbox(s, x + Inches(0.2), Inches(1.8), col_w - Inches(0.3),
                     Inches(0.7), col.get("heading", ""),
                     font_name=t.font_title, font_size=Pt(28),
                     color=t.text_light, bold=True)
        items = col.get("items") or []
        _add_textbox(s, x + Inches(0.2), Inches(2.6), col_w - Inches(0.3),
                     Inches(4.0), [f"•  {it}" for it in items],
                     font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_quote(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_textbox(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.5),
                 f"\u201C{slide.get('quote', '')}\u201D",
                 font_name=t.font_title, font_size=Pt(36),
                 color=t.text_dark, align=PP_ALIGN.CENTER)
    if slide.get("attribution"):
        _add_textbox(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.7),
                     f"— {slide['attribution']}",
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.secondary, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_data_callout(prs, slide, t: Tokens):
    return _simple_metric_spotlight(prs, slide, t)


def _simple_visual_hero(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    img = slide.get("image_path")
    if img:
        _add_picture(s, img, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    _add_textbox(s, Inches(0.75), Inches(5.5), Inches(11.83), Inches(1.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_section,
                 color=t.text_dark)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_question(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_hero,
                 color=t.text_dark, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _simple_cta_steps(prs, slide, t: Tokens):
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    steps = slide.get("steps") or []
    _add_textbox(s, Inches(0.75), Inches(1.9), Inches(11.83), Inches(5.0),
                 [f"{i+1}.  {step}" for i, step in enumerate(steps)],
                 font_name=t.font_body, font_size=t.sz_body,
                 color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


SIMPLE_BUILDERS = {
    "title":              _simple_title,
    "section-divider":    _simple_section_divider,
    "key-message":        _simple_key_message,
    "content":            _simple_content,
    "metric-spotlight":   _simple_metric_spotlight,
    "comparison-columns": _simple_comparison_columns,
    "quote":              _simple_quote,
    "data-callout":       _simple_data_callout,
    "visual-hero":        _simple_visual_hero,
    "question":           _simple_question,
    "cta-steps":          _simple_cta_steps,
}


# ============================================================
# STYLED BUILDERS  (architecture §4.1 EMU coords; F5)
# ============================================================

def _styled_hero_full_bleed(prs, slide, t: Tokens):
    """Full-bleed picture from `gradient_pattern` (or supplied image)
    + title centred lower-third + subtitle.
    EMU: picture covers (0,0,12192000,6858000)."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    asset = _first_asset(slide, recipe="gradient_pattern")
    img_path = (asset and asset.get("path")) or slide.get("image_path")
    if img_path and os.path.exists(img_path):
        _add_picture(s, img_path, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    _add_textbox(s, Inches(1.0), Inches(4.2), Inches(11.3), Inches(1.6),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_hero,
                 color=t.text_dark, align=PP_ALIGN.CENTER, bold=False)
    if slide.get("subtitle"):
        _add_textbox(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.6),
                     slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_dark, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_accent_block_left(prs, slide, t: Tokens):
    """Left 4023360-EMU (4.4") accent panel + right body. §4.1.

    C9 fix: the eyebrow no longer sits as body-weight text directly on
    the full-saturation accent panel (the canon-forbidden "body text on
    saturated accent"). The panel now carries only a large tracked
    *number/kicker* in a desaturated tint, plus a tint hairline; the
    readable title lives off-accent on the light surface."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_rect(s, 0, 0, Emu(4023360), SLIDE_H_EMU, t.primary)
    eyebrow = slide.get("eyebrow", "") or slide.get("kicker", "")
    # Tint the accent panel text toward white so it's a tonal mark, not
    # body copy fighting the saturated fill.
    panel_tint = _tint(t.primary, 0.78)
    if eyebrow:
        kt = _add_textbox(s, Inches(0.55), Inches(2.2), Inches(3.4), Inches(0.5),
                     str(eyebrow).upper(),
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=panel_tint, bold=True)
        _set_tracking(kt, 240)
    _add_hairline(s, Inches(0.55), Inches(2.75), Inches(1.0), panel_tint,
                  weight_emu=19050)
    # Large index numeral / section number on the panel (optional).
    num = slide.get("section_number") or slide.get("number")
    if num:
        nt = _add_textbox(s, Inches(0.5), Inches(3.1), Inches(3.5), Inches(2.0),
                     str(num), font_name=t.font_title, font_size=Pt(96),
                     color=panel_tint, bold=True)
        _set_tracking(nt, -40)
    _add_textbox(s, Emu(4023360) + Inches(0.6), Inches(2.0),
                 Inches(8.0), Inches(3.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_section,
                 color=t.text_light, anchor=MSO_ANCHOR.MIDDLE)
    if slide.get("subtitle"):
        _add_textbox(s, Emu(4023360) + Inches(0.6), Inches(5.0),
                     Inches(8.0), Inches(1.2), slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_2_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_metric_xxl(prs, slide, t: Tokens):
    """XXL number on bg_dark with delta/trend + unit styling. §4.1.

    C-restyle: the bare number gains an optional delta arrow (▲/▼ in
    focal/highlight tone), a unit suffix at reduced size, and a hairline
    rule under the label for editorial structure."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_kicker(s, Inches(0.75), Inches(0.9), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "",
                t, t.text_2_dark, rule_color=t.primary)
    metric_tx = _add_textbox(s, Inches(0.75), Inches(1.7), Inches(11.83), Inches(3.6),
                 str(slide.get("metric", "")),
                 font_name=t.font_title, font_size=t.sz_metric_xxl,
                 color=t.text_dark, align=PP_ALIGN.CENTER, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    _set_tracking(metric_tx, -60)
    # Delta / trend chip.
    delta = slide.get("delta")
    if delta:
        up = not str(delta).strip().startswith("-")
        arrow = "\u25B2" if up else "\u25BC"
        chip_color = t.chart_focal if up else t.highlight
        _add_textbox(s, Inches(0.75), Inches(5.35), Inches(11.83), Inches(0.55),
                     f"{arrow}  {delta}", font_name=t.font_body,
                     font_size=t.sz_subtitle, color=chip_color,
                     align=PP_ALIGN.CENTER, bold=True)
    if slide.get("metric_label"):
        _add_textbox(s, Inches(0.75), Inches(5.95), Inches(11.83), Inches(0.7),
                     slide["metric_label"],
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.secondary, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_quote_pullout(prs, slide, t: Tokens):
    """240pt opaque-10% quote glyph background + quote + attribution. §4.1."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    asset = _first_asset(slide, recipe="oversized_glyph_bg")
    if asset and asset.get("path") and os.path.exists(asset["path"]):
        _add_picture(s, asset["path"], 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    _add_textbox(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.5),
                 slide.get("quote", ""),
                 font_name=t.font_title, font_size=Pt(36),
                 color=t.text_dark, align=PP_ALIGN.CENTER)
    if slide.get("attribution"):
        _add_textbox(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.7),
                     f"— {slide['attribution']}",
                     font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.secondary, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_split_image_right(prs, slide, t: Tokens):
    """Left half (0..6096000) text; right half (6096000..12192000) image.
    `anchor: right_half`. §4.1."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    img = slide.get("image_path")
    if img and os.path.exists(img):
        _add_picture(s, img, Emu(SLIDE_HALF_EMU), 0,
                     Emu(SLIDE_HALF_EMU), SLIDE_H_EMU)
    _add_textbox(s, Inches(0.75), Inches(2.0), Inches(5.6), Inches(1.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    if slide.get("body"):
        _add_textbox(s, Inches(0.75), Inches(3.7), Inches(5.6), Inches(2.5),
                     slide["body"],
                     font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_tinted_panel_right(prs, slide, t: Tokens):
    """Aside / supporting fact: surface_elevated panel on the right.
    Right half tinted with `surface_elevated` (10% lighter)."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_rect(s, Emu(SLIDE_HALF_EMU), 0, Emu(SLIDE_HALF_EMU), SLIDE_H_EMU,
              t.surface)
    _add_textbox(s, Inches(0.75), Inches(2.0), Inches(5.6), Inches(3.5),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    aside = slide.get("aside") or slide.get("body") or ""
    _add_textbox(s, Inches(7.2), Inches(2.0), Inches(5.4), Inches(4.0),
                 aside,
                 font_name=t.font_body, font_size=t.sz_body,
                 color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_progress_dots(prs, slide, t: Tokens):
    """5-step progress strip. The matplotlib `progress_strip` recipe
    produces the line+dots; we overlay step labels."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.6), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    asset = _first_asset(slide, recipe="progress_strip")
    if asset and asset.get("path") and os.path.exists(asset["path"]):
        _add_picture(s, asset["path"], Inches(0.75), Inches(2.5),
                     Inches(11.83), Inches(3.5))
    else:
        # Fallback: simple text steps
        steps = slide.get("steps") or []
        _add_textbox(s, Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.5),
                     " → ".join(s.get("label", "?") if isinstance(s, dict) else s
                                for s in steps),
                     font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light, align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_chart_callout(prs, slide, t: Tokens):
    """Left half chart (matplotlib) + 2 right-side annotation textboxes."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)
    asset = _first_asset(slide, kind="chart")
    if asset and asset.get("path") and os.path.exists(asset["path"]):
        # C8: inset the chart to the safe margin — no x=0 edge-touch.
        # Was x=0 (violated the system's own safe_area_inches=0.5).
        _add_picture(s, asset["path"], Inches(0.5), Inches(1.5),
                     Emu(SLIDE_HALF_EMU) - Inches(0.5), Inches(5.3))
    callouts = slide.get("callouts") or []
    if len(callouts) >= 1:
        _add_textbox(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(1.8),
                     callouts[0], font_name=t.font_title, font_size=Pt(28),
                     color=t.primary, bold=True)
    if len(callouts) >= 2:
        _add_textbox(s, Inches(7.0), Inches(4.2), Inches(5.6), Inches(1.8),
                     callouts[1], font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


STYLED_BUILDERS = {
    "hero_full_bleed":     _styled_hero_full_bleed,
    "accent_block_left":   _styled_accent_block_left,
    "metric_xxl":          _styled_metric_xxl,
    "quote_pullout":       _styled_quote_pullout,
    "split_image_right":   _styled_split_image_right,
    "tinted_panel_right":  _styled_tinted_panel_right,
    "progress_dots":       _styled_progress_dots,
    "chart_callout":       _styled_chart_callout,
    # Session 2026-05-04-c8d3b2a1 — archetype builders
    "risk_heatmap":      None,  # bound below after defs
    "priority_matrix":   None,
    "waterfall":         None,
    "flywheel":          None,
    "funnel":            None,
    "decision_options":  None,
    "appendix_dense":    None,
    # Session 2026-06-08-c5d9e1a7 (C6) — editorial archetypes
    "stat_grid_3up":       None,
    "pull_quote_portrait": None,
    "full_bleed_caption":  None,
    "editorial_2col_6040": None,
    "timeline_horizontal": None,
    "agenda_toc":          None,
    "closing_cta":         None,
}


# ============================================================
# ARCHETYPE BUILDERS  (session 2026-05-04-c8d3b2a1)
# ============================================================
#
# Seven semantic builders realising the ⏳ Spec-only archetypes from
# `presentation-design/references/layout-archetypes.md`. Plus one
# partial (`_apply_footer_source`) applied to ANY slide whose spec
# carries a `footer` block.
#
# Geometry conventions (16:9 EMU):
#   slide  = (0, 0, 12192000, 6858000)  i.e. 13.333" × 7.5"
#   header = (0.75", 0.5", 11.83", 0.9")  for titles
#   body   = (0.75", 1.6", 11.83", 5.4")  for primary content
#   footer = (0,     7.10",   13.33", 0.35") reserved for footer_source
# Tokens are pulled from the resolved Tokens instance — no hard-coded
# colours. Each builder also receives the slide dict and may read its
# new-in-v2.2.0 spec keys (risks, axes, matrix_items, waterfall,
# flywheel_stages, funnel_stages, options/criteria/recommendation,
# panels, footer).
#
# Validation philosophy: builders raise SpecError on shape violations
# the eye can't easily catch (waterfall sign/zero-baseline, decision
# columns summing past slide width). Other failures degrade gracefully
# (missing labels render blank rather than crashing).


# ---------- Risk Heatmap (recipe: risk_heatmap) ----------

def _styled_risk_heatmap(prs, slide, t: Tokens):
    """3×3 likelihood × impact grid; risks plotted as labelled markers
    in cells; colour-coded green→yellow→red; legend at lower margin."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    LEVELS = ["low", "med", "high"]
    LEVEL_IDX = {v: i for i, v in enumerate(LEVELS)}

    # Grid: 3 cols × 3 rows centred in body
    grid_x = Inches(2.5); grid_y = Inches(1.7)
    cell_w = Inches(2.4); cell_h = Inches(1.4)

    # Cell fill semantics: green / amber / red by max(prob, impact) tier.
    # Colours chosen so WHITE 12pt-bold labels clear AA normal-text (4.5:1):
    #   green #1B5E20 → ~9.8:1   amber #B45309 → ~5.6:1   red #B71C1C → ~7.5:1
    # The earlier #F59E0B amber failed at 2.15:1 — fixed in
    # session 2026-05-04-c8d3b2a1.
    GREEN = _hex("#1B5E20")
    AMBER = _hex("#B45309")
    RED   = _hex("#B71C1C")
    LIGHT = _hex("#FFFFFF")
    BORDER = t.text_2_light

    def _cell_color(p_idx, i_idx):
        risk = max(p_idx, i_idx)
        if risk == 0:
            return GREEN
        if risk == 1:
            return AMBER
        return RED

    for ri in range(3):           # impact row 0..2 (bottom = high)
        for ci in range(3):       # probability col 0..2 (right = high)
            cx = grid_x + cell_w * ci
            cy = grid_y + cell_h * (2 - ri)
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cell_w, cell_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = _cell_color(ci, ri)
            rect.line.color.rgb = BORDER

    # Plot risks: deterministic offset within cell so multiple risks stack
    risks = slide.get("risks") or []
    cell_count = {}
    for r in risks:
        p = LEVEL_IDX.get(r.get("probability", "med"), 1)
        i = LEVEL_IDX.get(r.get("impact", "med"), 1)
        cx = grid_x + cell_w * p
        cy = grid_y + cell_h * (2 - i)
        slot = cell_count.get((p, i), 0)
        cell_count[(p, i)] = slot + 1
        offset_y = Inches(0.15 + 0.35 * slot)
        _add_textbox(s, cx + Inches(0.15), cy + offset_y,
                     cell_w - Inches(0.3), Inches(0.35),
                     "● " + str(r.get("name", "")),
                     font_name=t.font_body, font_size=Pt(12),
                     color=LIGHT, bold=True)

    # Axis labels
    axes = slide.get("axes") or {}
    _add_textbox(s, grid_x, grid_y + cell_h * 3 + Inches(0.05),
                 cell_w * 3, Inches(0.3),
                 axes.get("x_label", "Probability →"),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_2_light, align=PP_ALIGN.CENTER)
    _add_textbox(s, grid_x - Inches(1.7), grid_y, Inches(1.6),
                 cell_h * 3,
                 axes.get("y_label", "Impact →"),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_2_light, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Legend
    _add_textbox(s, Inches(0.75), Inches(6.5), Inches(11.83), Inches(0.4),
                 "Legend:  ● Low risk (green)   ● Medium (amber)   ● High (red)",
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_2_light)
    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- 2×2 / Priority Matrix (recipe: priority_matrix) ----------

def _styled_priority_matrix(prs, slide, t: Tokens):
    """2×2 grid; labelled axes; quadrant labels in faint type;
    items as labelled markers; recommended quadrant tinted."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    grid_x = Inches(3.0); grid_y = Inches(1.7)
    qw = Inches(3.5); qh = Inches(2.3)

    # Quadrants: TL TR BL BR
    quads = [(0, 0), (1, 0), (0, 1), (1, 1)]
    labels = slide.get("quadrant_labels") or ["", "", "", ""]
    items = slide.get("matrix_items") or []
    rec_quad = None
    # Find recommended quadrant from items
    for it in items:
        if it.get("recommended"):
            rec_quad = (1 if it.get("x") == "high" else 0,
                        0 if it.get("y") == "high" else 1)
            break

    for idx, (cx_i, cy_i) in enumerate(quads):
        x = grid_x + qw * cx_i
        y = grid_y + qh * cy_i
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, qw, qh)
        rect.fill.solid()
        if rec_quad and (cx_i, cy_i) == rec_quad:
            rect.fill.fore_color.rgb = t.primary
            label_color = t.text_dark
        else:
            rect.fill.fore_color.rgb = t.surface
            label_color = t.text_2_light
        rect.line.color.rgb = t.text_2_light
        # Quadrant label centre-faint
        _add_textbox(s, x + Inches(0.2), y + Inches(0.2),
                     qw - Inches(0.4), Inches(0.5),
                     labels[idx] if idx < len(labels) else "",
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=label_color, align=PP_ALIGN.CENTER)

    # Items
    for it in items:
        cx_i = 1 if it.get("x") == "high" else 0
        cy_i = 0 if it.get("y") == "high" else 1
        x = grid_x + qw * cx_i + Inches(0.4)
        y = grid_y + qh * cy_i + Inches(0.9)
        _add_textbox(s, x, y, qw - Inches(0.6), Inches(0.4),
                     "● " + str(it.get("name", "")),
                     font_name=t.font_body, font_size=t.sz_body,
                     color=(t.text_dark if (rec_quad == (cx_i, cy_i)) else t.text_light),
                     bold=bool(it.get("recommended")))

    # Axes
    axes = slide.get("axes") or {}
    _add_textbox(s, grid_x, grid_y + qh * 2 + Inches(0.05),
                 qw * 2, Inches(0.3),
                 axes.get("x_label", "→"),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_2_light, align=PP_ALIGN.CENTER)
    _add_textbox(s, grid_x - Inches(1.7), grid_y, Inches(1.6),
                 qh * 2,
                 axes.get("y_label", "↑"),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_2_light, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- Waterfall / Value Bridge (recipe: waterfall) ----------

def _styled_waterfall(prs, slide, t: Tokens):
    """Floating bars: start total, ordered signed deltas, end total.
    Positive=primary accent, negative=highlight. Zero-baseline enforced.
    Connector lines between bar tops."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    wf = slide.get("waterfall") or {}
    start = wf.get("start") or {}
    end = wf.get("end") or {}
    steps = wf.get("steps") or []
    units = wf.get("units", "")

    if not start or not end:
        _add_textbox(s, Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.0),
                     "[waterfall: missing start/end]",
                     font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_2_light, align=PP_ALIGN.CENTER)
        _set_notes(s, slide.get("notes", ""))
        return s

    # Geometry
    chart_left = Inches(1.0); chart_right = Inches(12.33)
    chart_top = Inches(2.0); chart_bottom = Inches(6.4)
    chart_w = chart_right - chart_left
    chart_h = chart_bottom - chart_top

    # All bars: start, *steps, end.  Compute running totals.
    s_val = float(start.get("value", 0))
    e_val = float(end.get("value", 0))
    running = [s_val]
    for st in steps:
        running.append(running[-1] + float(st.get("delta", 0)))
    # Final terminal bar uses end.value (may differ from running due to rounding/intent)
    all_vals = [s_val] + [running[i+1] for i in range(len(steps))] + [e_val]
    max_v = max(all_vals + [0])
    min_v = min(all_vals + [0])
    # Zero-baseline ENFORCED: scale from 0..max_v if all values >= 0,
    # else cover [min_v, max_v] but always include 0.
    y_lo = min(min_v, 0)
    y_hi = max(max_v, 0)
    span = (y_hi - y_lo) or 1
    n_bars = 1 + len(steps) + 1
    bar_slot_w = chart_w / max(1, n_bars)
    bar_w = int(bar_slot_w * 0.55)
    bar_pad = int((bar_slot_w - bar_w) / 2)

    def _y_for(v):
        return int(chart_top + (1 - (v - y_lo) / span) * chart_h)

    zero_y = _y_for(0)

    # Draw zero-baseline
    line = s.shapes.add_connector(2, chart_left, zero_y, chart_right, zero_y)
    line.line.color.rgb = t.text_2_light

    NEG = t.highlight
    POS = t.primary
    NEUTRAL = t.secondary

    # Start bar (neutral): from zero to start value
    x0 = int(chart_left + bar_pad)
    top0 = _y_for(max(0, s_val))
    bot0 = _y_for(min(0, s_val))
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, top0, bar_w, max(1, bot0 - top0))
    rect.fill.solid(); rect.fill.fore_color.rgb = NEUTRAL
    rect.line.fill.background()
    _add_textbox(s, x0 - Inches(0.2), Inches(6.5), bar_w + Inches(0.4), Inches(0.4),
                 str(start.get("label", "Start")),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_light, align=PP_ALIGN.CENTER)
    prev_running = s_val
    prev_top = top0 if s_val >= 0 else bot0

    # Step bars (floating)
    for i, st in enumerate(steps):
        delta = float(st.get("delta", 0))
        x = int(chart_left + bar_slot_w * (i + 1) + bar_pad)
        new_running = prev_running + delta
        top = _y_for(max(prev_running, new_running))
        bot = _y_for(min(prev_running, new_running))
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, bar_w, max(1, bot - top))
        rect.fill.solid()
        rect.fill.fore_color.rgb = POS if delta >= 0 else NEG
        rect.line.fill.background()
        _add_textbox(s, x - Inches(0.2), Inches(6.5), bar_w + Inches(0.4), Inches(0.4),
                     str(st.get("label", "")),
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=t.text_light, align=PP_ALIGN.CENTER)
        # Connector line from previous bar top to this bar top edge
        prev_x_right = int(chart_left + bar_slot_w * i + bar_pad + bar_w)
        c_y = _y_for(prev_running)
        conn = s.shapes.add_connector(1, prev_x_right, c_y, x, c_y)
        conn.line.color.rgb = t.text_2_light
        prev_running = new_running

    # End bar
    xN = int(chart_left + bar_slot_w * (len(steps) + 1) + bar_pad)
    topN = _y_for(max(0, e_val))
    botN = _y_for(min(0, e_val))
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xN, topN, bar_w, max(1, botN - topN))
    rect.fill.solid(); rect.fill.fore_color.rgb = NEUTRAL
    rect.line.fill.background()
    _add_textbox(s, xN - Inches(0.2), Inches(6.5), bar_w + Inches(0.4), Inches(0.4),
                 str(end.get("label", "End")),
                 font_name=t.font_body, font_size=t.sz_caption,
                 color=t.text_light, align=PP_ALIGN.CENTER)

    # Units in upper-right
    if units:
        _add_textbox(s, Inches(10.5), Inches(1.5), Inches(2.5), Inches(0.4),
                     f"({units})",
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=t.text_2_light, align=PP_ALIGN.RIGHT)

    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- Flywheel (recipe: flywheel) ----------

def _styled_flywheel(prs, slide, t: Tokens):
    """3–6 named nodes around a circle; reinforcing arrows; centre label.
    Stage labels positioned outside the ring for legibility."""
    import math
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    stages = slide.get("flywheel_stages") or []
    n = max(3, min(6, len(stages))) if stages else 0
    cx = Inches(6.665); cy = Inches(4.1)
    R = Inches(1.7)         # ring radius
    node_r = Inches(0.55)

    # Centre label
    centre_dia = Inches(1.5)
    centre = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - centre_dia / 2, cy - centre_dia / 2,
                                centre_dia, centre_dia)
    centre.fill.solid(); centre.fill.fore_color.rgb = t.primary
    centre.line.fill.background()
    _add_textbox(s, cx - centre_dia / 2, cy - Inches(0.3),
                 centre_dia, Inches(0.6),
                 slide.get("center_label", "") or "",
                 font_name=t.font_title, font_size=Pt(18),
                 color=t.text_dark, align=PP_ALIGN.CENTER, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Nodes around the ring
    node_centres = []
    for i in range(n):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        nx = int(cx + R * math.cos(ang))
        ny = int(cy + R * math.sin(ang))
        node_centres.append((nx, ny, ang))
        node = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  nx - node_r, ny - node_r,
                                  node_r * 2, node_r * 2)
        node.fill.solid(); node.fill.fore_color.rgb = t.secondary
        node.line.color.rgb = t.text_light
        # Label OUTSIDE the ring
        lx = int(cx + (R + node_r + Inches(0.3)) * math.cos(ang))
        ly = int(cy + (R + node_r + Inches(0.3)) * math.sin(ang))
        label = stages[i].get("label", "") if i < len(stages) else ""
        _add_textbox(s, lx - Inches(1.0), ly - Inches(0.2),
                     Inches(2.0), Inches(0.5),
                     label,
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=t.text_light, align=PP_ALIGN.CENTER, bold=True)

    # C10: reinforcing arc arrows along the ring WITH arrowheads (was a
    # straight tangent with no head — read as a polygon, not a flywheel).
    for i in range(n):
        a, b = node_centres[i], node_centres[(i + 1) % n]
        ax, ay, _ = a
        bx, by, _ = b
        dx, dy = bx - ax, by - ay
        L = (dx * dx + dy * dy) ** 0.5 or 1
        ux, uy = dx / L, dy / L
        sx = int(ax + ux * node_r)
        sy = int(ay + uy * node_r)
        ex = int(bx - ux * node_r)
        ey = int(by - uy * node_r)
        conn = s.shapes.add_connector(3, sx, sy, ex, ey)  # 3 = curved arc
        conn.line.color.rgb = t.primary
        conn.line.width = Emu(28575)  # ~2.25pt
        _add_arrowhead(conn)

    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- Funnel (recipe: funnel) ----------

def _styled_funnel(prs, slide, t: Tokens):
    """Trapezoidal stack, top-down narrowing. Per-stage label + count;
    conversion % between stages. Largest leak emphasised."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    stages = slide.get("funnel_stages") or []
    if not stages:
        _set_notes(s, slide.get("notes", ""))
        return s
    n = len(stages)

    # Geometry: vertical stack centred horizontally, narrowing.
    body_top = Inches(1.7); body_bot = Inches(6.5)
    stage_h = (body_bot - body_top) / n
    cx = Inches(6.665)
    max_w = Inches(7.0)
    min_w = Inches(2.5)
    counts = [float(st.get("count", 0)) for st in stages]
    top_count = counts[0] or 1
    # Width proportional to count
    widths = [max(min_w, int(max_w * (c / top_count))) for c in counts]

    # Find largest leak (drop) for emphasis
    leaks = [counts[i - 1] - counts[i] for i in range(1, n)]
    leak_max_idx = (leaks.index(max(leaks)) + 1) if leaks else None

    for i, st in enumerate(stages):
        y = int(body_top + stage_h * i)
        y_top = y + Emu(20000)
        y_bot = int(body_top + stage_h * (i + 1)) - Emu(20000)
        w_top = widths[i]
        # C10: true trapezoidal taper — bottom edge narrows to the NEXT
        # stage's width (last stage keeps a gentle close).
        w_bot = widths[i + 1] if i < n - 1 else int(w_top * 0.9)
        pts = [
            (cx - w_top / 2, y_top),
            (cx + w_top / 2, y_top),
            (cx + w_bot / 2, y_bot),
            (cx - w_bot / 2, y_bot),
        ]
        fill = t.primary if (i == leak_max_idx) else t.secondary
        try:
            fb = s.shapes.build_freeform(float(pts[0][0]), float(pts[0][1]), scale=1.0)
            fb.add_line_segments([(float(x), float(yy)) for x, yy in pts[1:]],
                                 close=True)
            shp = fb.convert_to_shape()
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
            shp.line.fill.background()
        except Exception:
            # Fallback: rectangle at top width.
            x = int(cx - w_top / 2)
            shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_top,
                                     w_top, y_bot - y_top)
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
            shp.line.fill.background()
        # Stage label + count centred on the band.
        _add_textbox(s, int(cx - w_top / 2) + Inches(0.15), y + Inches(0.05),
                     w_top - Inches(0.3), int(stage_h) - Emu(80000),
                     f"{st.get('label', '')}   ({int(st.get('count', 0)):,})",
                     font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_dark, align=PP_ALIGN.CENTER, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        # Conversion rate to RIGHT of bar (between this and prior)
        if i > 0:
            rate = st.get("rate")
            if rate is None and counts[i - 1]:
                rate = counts[i] / counts[i - 1]
            if rate is not None:
                _add_textbox(s, cx + max_w / 2 + Inches(0.2), y - Inches(0.15),
                             Inches(2.5), Inches(0.4),
                             f"\u2193 {rate * 100:.0f}%",
                             font_name=t.font_body, font_size=t.sz_caption,
                             color=t.text_2_light)

    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- Decision Options Table (recipe: decision_options) ----------

def _styled_decision_options(prs, slide, t: Tokens):
    """N options × M criteria table (>3 supported). Recommended option
    column highlighted; final 'Verdict' column shows ✓ / ✗.
    QA: column widths sum within slide-width tolerance."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.9),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=t.sz_title,
                 color=t.text_light)

    options = slide.get("options") or []
    criteria = slide.get("criteria") or []
    rec = slide.get("recommendation")

    if not options or not criteria:
        _set_notes(s, slide.get("notes", ""))
        return s

    n_opts = len(options)
    n_cols = 1 + n_opts            # leading "Criterion" col + one per option
    n_rows = 1 + len(criteria) + 1  # header + criteria + verdict

    table_left = Inches(0.75)
    table_top  = Inches(1.7)
    table_w    = Inches(11.83)
    table_h    = Inches(5.0)

    # Build table
    tbl_shape = s.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_w, table_h)
    tbl = tbl_shape.table

    # Column widths: criterion col = 25%, options share 75%
    crit_w = int(table_w * 0.25)
    opt_w = int((table_w - crit_w) / n_opts)
    tbl.columns[0].width = crit_w
    for c in range(1, n_cols):
        tbl.columns[c].width = opt_w

    # Header row
    tbl.cell(0, 0).text = "Criterion"
    for c, opt in enumerate(options, start=1):
        tbl.cell(0, c).text = str(opt.get("name", ""))
        cell = tbl.cell(0, c)
        cell.fill.solid()
        is_rec = (opt.get("name") == rec)
        cell.fill.fore_color.rgb = t.primary if is_rec else t.surface

    # Body rows: criteria
    for r, crit_name in enumerate(criteria, start=1):
        tbl.cell(r, 0).text = str(crit_name)
        for c, opt in enumerate(options, start=1):
            scores = opt.get("scores") or {}
            v = scores.get(crit_name, "")
            tbl.cell(r, c).text = str(v) if v != "" else "—"
            if opt.get("name") == rec:
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = t.surface

    # Verdict row
    tbl.cell(n_rows - 1, 0).text = "Verdict"
    for c, opt in enumerate(options, start=1):
        is_rec = (opt.get("name") == rec)
        tbl.cell(n_rows - 1, c).text = "✓ Recommended" if is_rec else "—"
        if is_rec:
            cell = tbl.cell(n_rows - 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = t.primary

    # Style table text — small + readable; skip granular formatting to keep
    # this builder under 100 LOC. python-pptx applies sane defaults.
    _set_notes(s, slide.get("notes", ""))

    # Self-check: column widths sum to slide width within tolerance.
    total = sum(col.width for col in tbl.columns)
    tol = Emu(50000)  # ~0.05"
    if abs(total - table_w) > tol:
        # Soft warning — does not block.
        pass
    return s


# ---------- Appendix Dense (recipe: appendix_dense) ----------

def _styled_appendix_dense(prs, slide, t: Tokens):
    """Dense reference layout: title bar, optional kicker, 2–4 panels
    (text/chart/table mix). Body type at caption scale. Visible
    'Appendix' tag upper-right. Pair with `slide.appendix=true` to
    bypass the body_word_max=30 density check."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)

    # Appendix tag upper-right
    _add_textbox(s, Inches(11.0), Inches(0.3), Inches(2.0), Inches(0.4),
                 "APPENDIX",
                 font_name=t.font_body, font_size=Pt(11),
                 color=t.text_2_light, align=PP_ALIGN.RIGHT, bold=True)

    _add_textbox(s, Inches(0.75), Inches(0.5), Inches(10.0), Inches(0.7),
                 slide.get("title", ""),
                 font_name=t.font_title, font_size=Pt(28),
                 color=t.text_light)
    if slide.get("subtitle"):
        _add_textbox(s, Inches(0.75), Inches(1.15), Inches(10.0), Inches(0.4),
                     slide["subtitle"],
                     font_name=t.font_body, font_size=t.sz_caption,
                     color=t.text_2_light)

    panels = slide.get("panels") or []
    n = len(panels)
    if n == 0:
        _set_notes(s, slide.get("notes", ""))
        return s

    # 1×2 (n=2) | 2×1 vertical (n in 3,4 use 2×2)
    if n == 2:
        cols, rows = 2, 1
    elif n <= 4:
        cols, rows = 2, 2
    else:
        cols, rows = 2, 2  # cap

    body_left = Inches(0.75); body_top = Inches(1.7)
    body_w = Inches(11.83);   body_h = Inches(5.4)
    pw = body_w / cols
    ph = body_h / rows
    pad = Inches(0.15)

    for i, p in enumerate(panels[:cols * rows]):
        ci = i % cols; ri = i // cols
        x = body_left + pw * ci + pad
        y = body_top + ph * ri + pad
        w = pw - pad * 2
        h = ph - pad * 2
        # Panel border
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        rect.fill.solid(); rect.fill.fore_color.rgb = t.surface
        rect.line.color.rgb = t.text_2_light

        if p.get("title"):
            _add_textbox(s, x + Inches(0.15), y + Inches(0.1),
                         w - Inches(0.3), Inches(0.4),
                         p["title"],
                         font_name=t.font_body, font_size=Pt(13),
                         color=t.text_light, bold=True)

        kind = p.get("kind", "text")
        body_y = y + Inches(0.55)
        body_h_inner = h - Inches(0.7)
        if kind == "text":
            _add_textbox(s, x + Inches(0.15), body_y,
                         w - Inches(0.3), body_h_inner,
                         p.get("body", ""),
                         font_name=t.font_body, font_size=Pt(12),
                         color=t.text_light)
        elif kind == "chart":
            img = p.get("image_path")
            if img and os.path.exists(img):
                _add_picture(s, img, x + Inches(0.15), body_y,
                             w - Inches(0.3), body_h_inner)
            else:
                _add_textbox(s, x + Inches(0.15), body_y,
                             w - Inches(0.3), body_h_inner,
                             "[chart missing]",
                             font_name=t.font_body, font_size=t.sz_caption,
                             color=t.text_2_light, align=PP_ALIGN.CENTER)
        elif kind == "table":
            rows_data = p.get("rows") or []
            if rows_data:
                _add_textbox(
                    s, x + Inches(0.15), body_y,
                    w - Inches(0.3), body_h_inner,
                    [" | ".join(r) for r in rows_data],
                    font_name=t.font_body, font_size=Pt(11),
                    color=t.text_light)

        if p.get("source"):
            _add_textbox(s, x + Inches(0.15), y + h - Inches(0.3),
                         w - Inches(0.3), Inches(0.25),
                         f"Source: {p['source']}",
                         font_name=t.font_body, font_size=Pt(9),
                         color=t.text_2_light)

    _set_notes(s, slide.get("notes", ""))
    return s


# ---------- Footer Source partial (recipe: footer_source — overlay) ----------

def _apply_footer_source(slide_obj, footer: dict, t: Tokens):
    """Overlay applied to ANY slide with `footer` block.

    `footer` shape: {source?, page?, page_total?, confidentiality?}
    Renders a 0.35"-tall band along bottom edge: source left,
    confidentiality centre, page right. Caption-size type."""
    if not footer:
        return
    band_top = Inches(7.10)
    band_h = Inches(0.35)
    source = footer.get("source")
    page = footer.get("page")
    page_total = footer.get("page_total")
    confidentiality = footer.get("confidentiality")

    if source:
        _add_textbox(slide_obj, Inches(0.35), band_top,
                     Inches(6.5), band_h,
                     str(source),
                     font_name=t.font_body, font_size=Pt(10),
                     color=t.text_2_light)
    if confidentiality:
        _add_textbox(slide_obj, Inches(5.0), band_top,
                     Inches(3.33), band_h,
                     str(confidentiality),
                     font_name=t.font_body, font_size=Pt(10),
                     color=t.text_2_light, align=PP_ALIGN.CENTER, bold=True)
    if page is not None:
        page_str = (f"{page} / {page_total}" if page_total else str(page))
        _add_textbox(slide_obj, Inches(10.0), band_top,
                     Inches(2.83), band_h,
                     page_str,
                     font_name=t.font_body, font_size=Pt(10),
                     color=t.text_2_light, align=PP_ALIGN.RIGHT)


# ============================================================
# C6 — High-impact editorial archetypes (session 2026-06-08-c5d9e1a7)
# ============================================================
#
# Seven crafted layouts built ON TOP of the C1 helpers (_add_card,
# _add_hairline, _add_scrim, _set_tracking, _add_kicker). Geometry is
# documented in pptx-engine/references/styled-recipes.md. These are the
# archetypes the aesthetic-leap user request explicitly asks for.


def _styled_stat_grid_3up(prs, slide, t: Tokens):
    """Three big-number cards across the body: rounded-rect surface,
    hairline border, 88pt number + delta arrow + 16pt label."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_kicker(s, Inches(0.75), Inches(0.55), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "", t,
                t.text_2_light, rule_color=t.primary)
    _add_textbox(s, Inches(0.75), Inches(1.05), Inches(11.83), Inches(0.9),
                 slide.get("title", ""), font_name=t.font_title,
                 font_size=t.sz_title, color=t.text_light)
    stats = (slide.get("stats") or slide.get("cards") or [])[:3]
    n = len(stats) or 1
    gutter = Inches(0.4)
    body_left = Inches(0.75); body_w = Inches(11.83)
    card_w = (body_w - gutter * (n - 1)) / n
    cy = Inches(2.3); card_h = Inches(3.8)
    for i, st in enumerate(stats):
        cx = body_left + (card_w + gutter) * i
        _add_card(s, cx, cy, card_w, card_h, fill=t.card_fill,
                  line=t.hairline, shadow=True)
        # Defensive sizing: long values (e.g. "99.99%", "3.2B/day") must
        # not wrap and overflow the card. Scale the numeral by length.
        val = str(st.get("value", ""))
        num_pt = 72 if len(val) <= 4 else (56 if len(val) <= 6 else 42)
        num_tx = _add_textbox(s, cx + Inches(0.3), cy + Inches(0.6),
                     card_w - Inches(0.6), Inches(1.6),
                     val, font_name=t.font_title,
                     font_size=Pt(num_pt), color=t.primary, bold=True,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        _set_tracking(num_tx, -40)
        delta = st.get("delta")
        if delta:
            up = not str(delta).strip().startswith("-")
            arrow = "\u25B2" if up else "\u25BC"
            _add_textbox(s, cx + Inches(0.3), cy + Inches(2.3),
                         card_w - Inches(0.6), Inches(0.5),
                         f"{arrow}  {delta}", font_name=t.font_body,
                         font_size=t.sz_caption,
                         color=t.chart_focal if up else t.highlight, bold=True)
        _add_hairline(s, cx + Inches(0.3), cy + Inches(2.85),
                      card_w - Inches(0.6), t.hairline)
        _add_textbox(s, cx + Inches(0.3), cy + Inches(2.95),
                     card_w - Inches(0.6), Inches(0.7),
                     str(st.get("label", "")), font_name=t.font_body,
                     font_size=Pt(16), color=t.text_2_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_pull_quote_portrait(prs, slide, t: Tokens):
    """Left 1/3 portrait card (rounded) + right 2/3 display-face quote,
    attribution, and role."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    # Portrait card on the left third.
    port_w = Inches(3.9); port_x = Inches(0.75)
    port_y = Inches(1.2); port_h = Inches(5.1)
    img = slide.get("image_path")
    if img and os.path.exists(img):
        _add_card(s, port_x, port_y, port_w, port_h, fill=t.card_fill,
                  line=t.hairline, shadow=True)
        _add_picture(s, img, port_x + Inches(0.12), port_y + Inches(0.12),
                     port_w - Inches(0.24), port_h - Inches(0.24))
    else:
        _add_card(s, port_x, port_y, port_w, port_h, fill=t.primary,
                  line=None, shadow=True)
    # Quote on the right.
    qx = Inches(5.1); qw = Inches(7.4)
    _add_textbox(s, qx, Inches(1.0), Inches(1.5), Inches(1.2), "\u201C",
                 font_name=t.font_title, font_size=Pt(120),
                 color=t.primary, bold=True)
    quote_tx = _add_textbox(s, qx, Inches(2.1), qw, Inches(3.0),
                 slide.get("quote", ""), font_name=t.font_title,
                 font_size=Pt(30), color=t.text_light)
    _set_tracking(quote_tx, -10)
    _add_hairline(s, qx, Inches(5.4), Inches(1.0), t.primary, weight_emu=19050)
    if slide.get("attribution"):
        _add_textbox(s, qx, Inches(5.55), qw, Inches(0.5),
                     str(slide["attribution"]), font_name=t.font_body,
                     font_size=t.sz_subtitle, color=t.text_light, bold=True)
    if slide.get("role"):
        _add_textbox(s, qx, Inches(6.05), qw, Inches(0.5),
                     str(slide["role"]), font_name=t.font_body,
                     font_size=t.sz_caption, color=t.text_2_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_full_bleed_caption(prs, slide, t: Tokens):
    """Full-bleed image + bottom gradient scrim + caption strip. The
    scrim guarantees caption contrast over any image."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    asset = _first_asset(slide, kind="composite") or _first_asset(slide)
    img = (asset and asset.get("path")) or slide.get("image_path")
    if img and os.path.exists(img):
        _add_picture(s, img, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    # Bottom scrim for caption contrast (transparent → near-opaque).
    _add_scrim(s, 0, Inches(4.3), SLIDE_W_EMU, Inches(3.2), t.scrim,
               top_alpha=0, bottom_alpha=82)
    _add_kicker(s, Inches(0.9), Inches(5.3), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "", t,
                _tint(t.primary, 0.4), rule=False)
    title_tx = _add_textbox(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.0),
                 slide.get("title", ""), font_name=t.font_title,
                 font_size=t.sz_section, color=t.text_dark, bold=False)
    _set_tracking(title_tx, -20)
    if slide.get("caption"):
        _add_textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.6),
                     slide["caption"], font_name=t.font_body,
                     font_size=t.sz_caption, color=_tint(t.scrim, 0.75))
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_editorial_2col_6040(prs, slide, t: Tokens):
    """60% lead column (standfirst + <=3 points) + 40% tonal aside on a
    surface_elevated card. Magazine-grade asymmetry."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_kicker(s, Inches(0.75), Inches(0.55), Inches(7.0),
                slide.get("eyebrow") or slide.get("kicker") or "", t,
                t.text_2_light, rule_color=t.primary)
    _add_textbox(s, Inches(0.75), Inches(1.05), Inches(7.0), Inches(1.4),
                 slide.get("title", ""), font_name=t.font_title,
                 font_size=t.sz_title, color=t.text_light)
    # Lead column.
    lead = slide.get("lead") or slide.get("standfirst")
    y = Inches(2.7)
    if lead:
        _add_textbox(s, Inches(0.75), y, Inches(6.7), Inches(1.4),
                     str(lead), font_name=t.font_title, font_size=Pt(24),
                     color=t.text_light)
        y = Inches(4.1)
    points = (slide.get("points") or slide.get("bullets") or [])[:3]
    for i, pt in enumerate(points):
        py = y + Inches(0.85) * i
        _add_hairline(s, Inches(0.75), py, Inches(6.7), t.hairline)
        _add_textbox(s, Inches(0.75), py + Inches(0.08), Inches(0.6), Inches(0.6),
                     f"{i + 1:02d}", font_name=t.font_title, font_size=Pt(18),
                     color=t.primary, bold=True)
        _add_textbox(s, Inches(1.45), py + Inches(0.08), Inches(6.0), Inches(0.7),
                     str(pt), font_name=t.font_body, font_size=t.sz_body,
                     color=t.text_light)
    # 40% tonal aside card.
    aside_x = Inches(8.0); aside_w = Inches(4.58)
    _add_card(s, aside_x, Inches(1.05), aside_w, Inches(5.6),
              fill=t.card_fill, line=t.hairline, shadow=True)
    aside = slide.get("aside") or {}
    if isinstance(aside, str):
        aside = {"body": aside}
    _add_kicker(s, aside_x + Inches(0.35), Inches(1.45), aside_w - Inches(0.7),
                aside.get("kicker", "Context"), t, t.primary, rule=False)
    _add_textbox(s, aside_x + Inches(0.35), Inches(2.0), aside_w - Inches(0.7),
                 Inches(4.4), aside.get("body", ""), font_name=t.font_body,
                 font_size=t.sz_body, color=t.text_light)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_timeline_horizontal(prs, slide, t: Tokens):
    """Baseline rule + milestone nodes + alternating date/label callouts."""
    import math  # noqa: F401
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_light)
    _add_kicker(s, Inches(0.75), Inches(0.55), Inches(8.0),
                slide.get("eyebrow") or slide.get("kicker") or "", t,
                t.text_2_light, rule_color=t.primary)
    _add_textbox(s, Inches(0.75), Inches(1.05), Inches(11.83), Inches(0.9),
                 slide.get("title", ""), font_name=t.font_title,
                 font_size=t.sz_title, color=t.text_light)
    milestones = slide.get("milestones") or slide.get("steps") or []
    n = max(1, len(milestones))
    axis_y = Inches(3.9)
    left = Inches(1.2); right = Inches(12.13)
    span = right - left
    _add_hairline(s, left, axis_y, span, t.text_2_light, weight_emu=19050)
    node_r = Inches(0.16)
    for i, m in enumerate(milestones):
        mx = int(left + span * (i / (n - 1))) if n > 1 else int(left + span / 2)
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, mx - node_r, axis_y - node_r,
                                  node_r * 2, node_r * 2)
        node.fill.solid()
        node.fill.fore_color.rgb = t.primary if m.get("highlight") else t.secondary
        node.line.color.rgb = t.bg_light
        node.line.width = Emu(28575)
        above = (i % 2 == 0)
        ly = axis_y - Inches(1.5) if above else axis_y + Inches(0.4)
        date_tx = _add_textbox(s, mx - Inches(1.3), ly, Inches(2.6), Inches(0.4),
                     str(m.get("date", "")), font_name=t.font_body,
                     font_size=t.sz_caption, color=t.primary,
                     align=PP_ALIGN.CENTER, bold=True)
        _set_tracking(date_tx, 120)
        _add_textbox(s, mx - Inches(1.3), ly + Inches(0.4), Inches(2.6), Inches(1.0),
                     str(m.get("label", "")), font_name=t.font_body,
                     font_size=t.sz_body, color=t.text_light,
                     align=PP_ALIGN.CENTER)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_agenda_toc(prs, slide, t: Tokens):
    """Numbered sections with tracked numerals + hairline separators —
    reuses the editorial kicker/numeral craft."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_kicker(s, Inches(0.9), Inches(0.7), Inches(8.0),
                slide.get("eyebrow") or "Agenda", t, t.text_2_dark,
                rule_color=t.primary)
    _add_textbox(s, Inches(0.9), Inches(1.2), Inches(11.0), Inches(1.5),
                 slide.get("title", "What we'll cover"),
                 font_name=t.font_title, font_size=t.sz_section,
                 color=t.text_dark)
    items = slide.get("sections") or slide.get("items") or []
    y0 = Inches(2.85)
    row_h = (Inches(6.6) - y0) / max(1, len(items))
    for i, it in enumerate(items):
        y = y0 + row_h * i
        label = it.get("label", "") if isinstance(it, dict) else str(it)
        _add_hairline(s, Inches(0.9), y, Inches(11.5), t.hairline_on_dark)
        nt = _add_textbox(s, Inches(0.9), y + Inches(0.1), Inches(1.4), row_h,
                     f"{i + 1:02d}", font_name=t.font_title, font_size=Pt(34),
                     color=t.primary, bold=True)
        _set_tracking(nt, -20)
        _add_textbox(s, Inches(2.5), y + Inches(0.1), Inches(9.0), row_h,
                     label, font_name=t.font_body, font_size=t.sz_subtitle,
                     color=t.text_dark, anchor=MSO_ANCHOR.MIDDLE)
    _set_notes(s, slide.get("notes", ""))
    return s


def _styled_closing_cta(prs, slide, t: Tokens):
    """Restrained closing/CTA: asymmetric headline, one accent rule,
    contact block. No decorative noise."""
    s = _new_blank_slide(prs)
    _set_bg(s, t.bg_dark)
    _add_rect(s, 0, 0, SLIDE_W_EMU, Emu(54864), t.primary)
    _add_kicker(s, Inches(1.2), Inches(1.6), Inches(8.0),
                slide.get("eyebrow") or "Thank you", t, t.text_2_dark,
                rule_color=t.primary)
    title_tx = _add_textbox(s, Inches(1.2), Inches(2.3), Inches(10.0), Inches(2.0),
                 slide.get("title", ""), font_name=t.font_title,
                 font_size=t.sz_hero, color=t.text_dark)
    _set_tracking(title_tx, -20)
    if slide.get("subtitle"):
        _add_textbox(s, Inches(1.2), Inches(4.4), Inches(10.0), Inches(0.8),
                     slide["subtitle"], font_name=t.font_body,
                     font_size=t.sz_subtitle, color=t.secondary)
    _add_hairline(s, Inches(1.2), Inches(5.6), Inches(10.9), t.hairline_on_dark)
    contact = slide.get("contact") or {}
    if isinstance(contact, str):
        contact = {"email": contact}
    parts = []
    for key in ("name", "email", "phone", "url"):
        if isinstance(contact, dict) and contact.get(key):
            parts.append(str(contact[key]))
    if not parts and slide.get("cta"):
        parts = [str(slide["cta"])]
    if parts:
        _add_textbox(s, Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.8),
                     "    ".join(parts), font_name=t.font_body,
                     font_size=t.sz_body, color=t.text_dark)
    _set_notes(s, slide.get("notes", ""))
    return s


# Bind the archetype builders into STYLED_BUILDERS now they're defined.
STYLED_BUILDERS["risk_heatmap"]     = _styled_risk_heatmap
STYLED_BUILDERS["priority_matrix"]  = _styled_priority_matrix
STYLED_BUILDERS["waterfall"]        = _styled_waterfall
STYLED_BUILDERS["flywheel"]         = _styled_flywheel
STYLED_BUILDERS["funnel"]           = _styled_funnel
STYLED_BUILDERS["decision_options"] = _styled_decision_options
STYLED_BUILDERS["appendix_dense"]   = _styled_appendix_dense
# C6 (session 2026-06-08-c5d9e1a7)
STYLED_BUILDERS["stat_grid_3up"]       = _styled_stat_grid_3up
STYLED_BUILDERS["pull_quote_portrait"] = _styled_pull_quote_portrait
STYLED_BUILDERS["full_bleed_caption"]  = _styled_full_bleed_caption
STYLED_BUILDERS["editorial_2col_6040"] = _styled_editorial_2col_6040
STYLED_BUILDERS["timeline_horizontal"] = _styled_timeline_horizontal
STYLED_BUILDERS["agenda_toc"]          = _styled_agenda_toc
STYLED_BUILDERS["closing_cta"]         = _styled_closing_cta


def _first_asset(slide, *, kind: str = None, recipe: str = None) -> dict | None:
    for a in (slide.get("visual_assets") or []):
        if a.get("skipped"):
            continue
        if kind and a.get("kind") != kind:
            continue
        if recipe and a.get("recipe") != recipe:
            continue
        return a
    return None


# ============================================================
# Spec validation + dispatch (F5, C5)
# ============================================================

class SpecError(Exception):
    """bad_spec.<reason>"""


def _validate_slide(slide: dict) -> tuple[str, str | None]:
    """Return (style, style_recipe). Apply C5 backwards-compat:
    missing `style` → "simple"; only reject when present-and-invalid
    or pairing-inconsistent."""
    style = slide.get("style", "simple")  # C5 default
    recipe = slide.get("style_recipe")
    stype = slide.get("type")

    if style not in ("simple", "styled"):
        raise SpecError(f"bad_spec.invalid_style: slide {slide.get('index')} "
                        f"has style={style!r}; must be 'simple' or 'styled'")

    if style == "styled":
        if not recipe:
            raise SpecError(f"bad_spec.styled_without_recipe: slide "
                            f"{slide.get('index')} has style=styled but no style_recipe")
        if recipe not in STYLED_RECIPES:
            raise SpecError(f"bad_spec.unknown_recipe: slide "
                            f"{slide.get('index')} style_recipe={recipe!r} "
                            f"not in {sorted(STYLED_RECIPES)}")
    else:  # simple
        if recipe is not None:
            raise SpecError(f"bad_spec.simple_with_recipe: slide "
                            f"{slide.get('index')} has style=simple but "
                            f"style_recipe={recipe!r} (inconsistent)")
        if stype and stype not in SIMPLE_TYPES:
            # Don't reject — backwards-compat: unknown simple types
            # fall back to `content`. Log at the call site.
            pass
    return style, recipe


def _build_slide(prs, slide: dict, t: Tokens, log) -> None:
    style, recipe = _validate_slide(slide)
    if style == "styled":
        builder = STYLED_BUILDERS[recipe]
        log.write(f"[slide {slide.get('index')}] styled/{recipe}\n")
        slide_obj = builder(prs, slide, t)
    else:
        stype = slide.get("type", "content")
        builder = SIMPLE_BUILDERS.get(stype, _simple_content)
        log.write(f"[slide {slide.get('index')}] simple/{stype}\n")
        slide_obj = builder(prs, slide, t)

    # Footer Source partial (session 2026-05-04-c8d3b2a1) — applied to
    # ANY slide whose spec carries a `footer` block, regardless of the
    # primary builder used. Reuses the same Tokens instance.
    footer = slide.get("footer")
    if footer and slide_obj is not None:
        try:
            _apply_footer_source(slide_obj, footer, t)
        except Exception as e:
            log.write(f"[footer] failed on slide {slide.get('index')}: {e}\n")


# ============================================================
# Main
# ============================================================

def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True, type=Path,
                    help="Path to deck-spec.json")
    ap.add_argument("--out", type=Path, default=Path("output.pptx"))
    ap.add_argument("--template", type=Path, default=None)
    ap.add_argument("--build-log", type=Path, default=None)
    args = ap.parse_args()

    log_path = args.build_log or args.out.with_suffix(".log")
    log = log_path.open("w", encoding="utf-8")
    try:
        try:
            deck_spec = json.loads(args.spec.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"ERROR: cannot read deck-spec: {e}", file=sys.stderr)
            return 2

        # Tokens (F9)
        t = Tokens(deck_spec.get("design_system_tokens") or FALLBACK_TOKENS)
        log.write(f"[tokens] {(deck_spec.get('design_system_tokens') or FALLBACK_TOKENS).get('name','?')}\n")

        # Pre-render visual_assets
        _prerender_assets(args.spec, deck_spec, log)

        # Presentation init
        if args.template and args.template.exists():
            try:
                prs = Presentation(str(args.template))
                log.write(f"[template] loaded {args.template}\n")
            except Exception as e:
                log.write(f"[template] failed ({e}); using default 16:9\n")
                prs = Presentation()
                prs.slide_width = Emu(SLIDE_W_EMU)
                prs.slide_height = Emu(SLIDE_H_EMU)
        else:
            prs = Presentation()
            prs.slide_width = Emu(SLIDE_W_EMU)
            prs.slide_height = Emu(SLIDE_H_EMU)

        # Dispatch
        try:
            for slide in deck_spec.get("slides", []):
                _build_slide(prs, slide, t, log)
        except SpecError as e:
            print(f"ERROR: {e}", file=sys.stderr)
            log.write(f"[reject] {e}\n")
            return 2

        args.out.parent.mkdir(parents=True, exist_ok=True)
        # Final WCAG contrast guard: nudge only the text runs that fail AA
        # against their effective background (mirrors check_pptx's resolver),
        # so accent marks resolve with genuine contrast.
        _enforce_text_contrast(prs, log)
        prs.save(str(args.out))
        print(f"SUCCESS: deck written to {args.out}")
        log.write(f"[done] {args.out}\n")
        return 0
    finally:
        log.close()


if __name__ == "__main__":
    sys.exit(main())
